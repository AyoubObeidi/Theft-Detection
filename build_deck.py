# -*- coding: utf-8 -*-
"""
Build a premium 30-slide pitch deck for the Theft Detection project.
Run:  py build_deck.py   ->  Theft_Detection_Pitch.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------- palette
NAVY_DEEP = RGBColor(0x0B, 0x18, 0x30)
GRAD2     = RGBColor(0x1B, 0x3A, 0x66)
NAVY      = RGBColor(0x0F, 0x1F, 0x38)
BLUE      = RGBColor(0x2E, 0x74, 0xB5)
BLUE_BR   = RGBColor(0x3B, 0x82, 0xF6)
ACCENT    = RGBColor(0x44, 0x72, 0xC4)
LIGHT     = RGBColor(0xDA, 0xE3, 0xF3)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT      = RGBColor(0xE8, 0xEE, 0xF7)
MUTED     = RGBColor(0x9D, 0xB2, 0xCE)
CARD      = RGBColor(0x14, 0x27, 0x47)
CARD_HI   = RGBColor(0x19, 0x33, 0x5C)
CARD_LINE = RGBColor(0x2A, 0x3F, 0x66)
GREEN     = RGBColor(0x28, 0xC8, 0x40)
RED       = RGBColor(0xFF, 0x5F, 0x57)
AMBER     = RGBColor(0xFE, 0xBC, 0x2E)

BODY = "Segoe UI"
TITLE = "Segoe UI Semibold"

LEFT, CENTER, RIGHT = PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT
TOP, MIDDLE, BOTTOM = MSO_ANCHOR.TOP, MSO_ANCHOR.MIDDLE, MSO_ANCHOR.BOTTOM

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.9)
CW = SW - 2 * MARGIN
SHOTS = "docs/shots"
LOGO = "docs/djelfa_logo.png"

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH


# ----------------------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def _noshadow(sh):
    sh.shadow.inherit = False
    return sh


def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = s.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    _noshadow(sh)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def solid_bg(s, color=NAVY):
    rect(s, 0, 0, SW, SH, fill=color)


def grad_bg(s, c1=NAVY_DEEP, c2=GRAD2, angle=60):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    sh.line.fill.background(); _noshadow(sh)
    try:
        sh.fill.gradient()
        stops = sh.fill.gradient_stops
        stops[0].position = 0.0; stops[0].color.rgb = c1
        stops[1].position = 1.0; stops[1].color.rgb = c2
        try:
            sh.fill.gradient_angle = angle
        except Exception:
            pass
    except Exception:
        sh.fill.solid(); sh.fill.fore_color.rgb = c1
    return sh


def tracking(run, pts):
    try:
        run.font._rPr.set('spc', str(int(pts * 100)))
    except Exception:
        pass


def text(s, l, t, w, h, lines, align=LEFT, anchor=TOP, track=None):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get('align', align)
        p._p.get_or_add_pPr().set('rtl', '0')  # force LTR (avoid Arabic-locale bidi reordering)
        if ln.get('sb') is not None:
            p.space_before = Pt(ln['sb'])
        if ln.get('sa') is not None:
            p.space_after = Pt(ln['sa'])
        if ln.get('ls') is not None:
            p.line_spacing = ln['ls']
        runs = ln['runs'] if 'runs' in ln else [ln]
        for rd in runs:
            r = p.add_run()
            r.text = '‪' + rd['t'] + '‬'  # LTR embedding: stop locale bidi reordering of $, /, digits
            r.font.size = Pt(rd.get('size', 18))
            r.font.bold = rd.get('bold', False)
            r.font.italic = rd.get('italic', False)
            r.font.name = rd.get('font', BODY)
            r.font.color.rgb = rd.get('color', TEXT)
            if rd.get('track'):
                tracking(r, rd['track'])
        if track:
            for r in p.runs:
                tracking(r, track)
    return tb


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def top_bar(s):
    rect(s, 0, 0, SW, Pt(6), fill=BLUE)


def footer(s, n):
    rect(s, MARGIN, SH - Inches(0.52), CW, Pt(0.75), fill=CARD_LINE)
    text(s, MARGIN, SH - Inches(0.47), Inches(5), Inches(0.3),
         [{'t': 'THEFT DETECTION', 'size': 8.5, 'color': MUTED, 'bold': True, 'track': 2}])
    text(s, SW - MARGIN - Inches(2), SH - Inches(0.47), Inches(2), Inches(0.3),
         [{'t': f'{n:02d} / 30', 'size': 8.5, 'color': MUTED, 'align': RIGHT}], align=RIGHT)


def header(s, kicker, title, n, tsize=32):
    solid_bg(s); top_bar(s); footer(s, n)
    if kicker:
        text(s, MARGIN, Inches(0.66), Inches(11), Inches(0.35),
             [{'t': kicker.upper(), 'size': 12, 'color': BLUE_BR, 'bold': True}], track=2.5)
    text(s, MARGIN, Inches(1.04), Inches(11.6), Inches(1.0),
         [{'t': title, 'size': tsize, 'color': WHITE, 'bold': True, 'font': TITLE}])


def card(s, l, t, w, h, fill=CARD, line=CARD_LINE, lw=1.0, radius=0.055):
    return rect(s, l, t, w, h, fill=fill, line=line, lw=lw,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)


def chip(s, l, t, sz, label, fill=BLUE, fg=WHITE, fsize=15):
    rect(s, l, t, sz, sz, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.28)
    text(s, l, t, sz, sz, [{'t': label, 'size': fsize, 'color': fg, 'bold': True, 'align': CENTER}],
         align=CENTER, anchor=MIDDLE)


def radar(s, cx, cy, color=BLUE):
    """decorative concentric rings (vision / scan motif)"""
    for r in [Inches(2.7), Inches(2.0), Inches(1.35), Inches(0.75)]:
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
        o.fill.background(); o.line.color.rgb = color; o.line.width = Pt(1.1); _noshadow(o)
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.12), cy - Inches(0.12), Inches(0.24), Inches(0.24))
    d.fill.solid(); d.fill.fore_color.rgb = BLUE_BR; d.line.fill.background(); _noshadow(d)


def shot(s, img, l, t, w, ar=1.6):
    """screenshot inside a browser-style frame"""
    bar = Inches(0.30)
    iw = w - Inches(0.08)
    ih = iw / ar
    card(s, l, t, w, bar + ih + Inches(0.06), fill=CARD_HI, line=CARD_LINE, radius=0.03)
    for i, c in enumerate([RED, AMBER, GREEN]):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.20) + i * Inches(0.24),
                               t + Inches(0.10), Inches(0.12), Inches(0.12))
        d.fill.solid(); d.fill.fore_color.rgb = c; d.line.fill.background(); _noshadow(d)
    s.shapes.add_picture(img, l + Inches(0.04), t + bar, width=iw)
    return t + bar + ih


def arrow(s, l, t, w, h, color=BLUE):
    a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); _noshadow(a)
    return a


# ============================================================================= SLIDE 1
s = slide(); grad_bg(s)
radar(s, Inches(11.3), Inches(3.6))
text(s, MARGIN, Inches(1.5), Inches(8), Inches(0.4),
     [{'t': 'AI VISION SECURITY  ·  RETAIL LOSS PREVENTION', 'size': 12, 'color': BLUE_BR, 'bold': True}], track=2.5)
text(s, MARGIN, Inches(2.0), Inches(9.2), Inches(2.6),
     [{'t': 'Theft Detection', 'size': 60, 'color': WHITE, 'bold': True, 'font': TITLE, 'sa': 6},
      {'t': 'AI that sees shoplifting before it happens.', 'size': 24, 'color': LIGHT}])
rect(s, MARGIN, Inches(4.55), Inches(0.8), Pt(3), fill=BLUE)
text(s, MARGIN, Inches(5.7), Inches(9), Inches(0.9),
     [{'t': 'Real-time pose & object intelligence for stores, supermarkets and smart facilities.',
       'size': 14, 'color': MUTED}])
text(s, MARGIN, Inches(6.5), Inches(11), Inches(0.4),
     [{'runs': [{'t': 'Live demo  ', 'size': 13, 'color': MUTED},
                {'t': 'theft-detection-dusky.vercel.app', 'size': 13, 'color': BLUE_BR, 'bold': True}]}])
try:
    s.shapes.add_picture(LOGO, SW - Inches(1.9), Inches(0.7), height=Inches(1.0))
except Exception:
    pass
notes(s, "LAYOUT: Title hero, left-aligned. ELEMENTS: oversized wordmark, blue rule, decorative radar rings (vision motif), university logo top-right. ANIMATION: Title Fade-up; radar rings Zoom/Grow on entrance; use Morph transition into Slide 2.")

# ============================================================================= SLIDE 2
s = slide(); grad_bg(s, angle=90); footer(s, 2)
text(s, 0, Inches(1.5), SW, Inches(0.5),
     [{'t': 'EVERY YEAR, RETAIL LOSES', 'size': 14, 'color': BLUE_BR, 'bold': True, 'align': CENTER}],
     align=CENTER, track=3)
text(s, 0, Inches(2.0), SW, Inches(2.6),
     [{'t': '$112B', 'size': 150, 'color': WHITE, 'bold': True, 'font': TITLE, 'align': CENTER}],
     align=CENTER, anchor=MIDDLE)
text(s, 0, Inches(4.7), SW, Inches(0.6),
     [{'t': 'to theft, fraud and shrink — most of it never seen in time.', 'size': 20, 'color': LIGHT, 'align': CENTER}],
     align=CENTER)
text(s, 0, Inches(6.4), SW, Inches(0.4),
     [{'t': '*Illustrative industry figure — replace with your sourced number.', 'size': 10, 'color': MUTED, 'align': CENTER}],
     align=CENTER)
notes(s, "LAYOUT: Single big-number statement, centered. ELEMENT: one hero numeral, supporting line. ANIMATION: numeral Zoom-in (Grow & Turn), supporting line Fade-up after. Keep on screen 3-4s for impact.")

# ============================================================================= SLIDE 3
s = slide(); header(s, "The Problem", "Theft is fast, silent — and human eyes miss it", 3)
items = [
    ("Seconds, not minutes", "A concealment gesture takes 2-3 seconds. Staff can't watch every aisle at once."),
    ("Walls of dead video", "Dozens of CCTV feeds, zero intelligence. Footage is reviewed only after the loss."),
    ("No proof, no response", "By the time anyone notices, there's no live alert and rarely usable evidence."),
]
y = Inches(2.45)
for i, (h, sub) in enumerate(items):
    rect(s, MARGIN, y + Inches(0.05), Pt(3), Inches(1.05), fill=BLUE)
    text(s, MARGIN + Inches(0.35), y, Inches(10.6), Inches(1.1),
         [{'t': h, 'size': 19, 'color': WHITE, 'bold': True, 'sa': 3},
          {'t': sub, 'size': 14, 'color': MUTED}])
    y += Inches(1.35)
notes(s, "LAYOUT: Stacked pain points with vertical accent rules. ANIMATION: each row Fade-up, staggered 0.15s. Speak to one pain per row.")

# ============================================================================= SLIDE 4
s = slide(); header(s, "Cost of Inaction", "What passive security really costs", 4)
stats = [("1 in 11", "shoppers shoplift — most never detected"),
         ("~2%", "of revenue lost to shrink, straight off margin"),
         ("0", "live evidence captured by legacy recording")]
gap = Inches(0.4); cwd = (CW - 2 * gap) / 3
for i, (num, lab) in enumerate(stats):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.5), cwd, Inches(3.0))
    text(s, x, Inches(2.95), cwd, Inches(1.2),
         [{'t': num, 'size': 52, 'color': BLUE_BR, 'bold': True, 'font': TITLE, 'align': CENTER}], align=CENTER)
    rect(s, x + cwd / 2 - Inches(0.4), Inches(4.25), Inches(0.8), Pt(2.5), fill=CARD_LINE)
    text(s, x + Inches(0.3), Inches(4.5), cwd - Inches(0.6), Inches(0.9),
         [{'t': lab, 'size': 14, 'color': LIGHT, 'align': CENTER}], align=CENTER)
notes(s, "LAYOUT: 3 stat cards. ANIMATION: cards Float-in (Up) staggered; numbers count-up feel via Zoom. *figures illustrative.")

# ============================================================================= SLIDE 5
s = slide(); header(s, "Why Now", "Legacy CCTV records the past. It can't prevent it.", 5, tsize=30)
gap = Inches(0.5); cwd = (CW - gap) / 2
old = ["Passive recording only", "Reviewed after the loss", "One guard, many screens", "No alerts, no analytics"]
new = ["Real-time AI analysis", "Alerts the moment it happens", "Watches every feed at once", "Evidence + insight, instantly"]
# left card (legacy)
card(s, MARGIN, Inches(2.4), cwd, Inches(3.7), fill=CARD)
text(s, MARGIN + Inches(0.4), Inches(2.7), cwd - Inches(0.8), Inches(0.5),
     [{'t': 'TRADITIONAL CCTV', 'size': 13, 'color': MUTED, 'bold': True}], track=2)
yy = Inches(3.35)
for it in old:
    text(s, MARGIN + Inches(0.4), yy, cwd - Inches(0.8), Inches(0.5),
         [{'runs': [{'t': '✕   ', 'size': 15, 'color': RED, 'bold': True}, {'t': it, 'size': 15, 'color': LIGHT}]}])
    yy += Inches(0.62)
# right card (ours)
x2 = MARGIN + cwd + gap
card(s, x2, Inches(2.4), cwd, Inches(3.7), fill=CARD_HI, line=BLUE, lw=1.5)
text(s, x2 + Inches(0.4), Inches(2.7), cwd - Inches(0.8), Inches(0.5),
     [{'t': 'THEFT DETECTION', 'size': 13, 'color': BLUE_BR, 'bold': True}], track=2)
yy = Inches(3.35)
for it in new:
    text(s, x2 + Inches(0.4), yy, cwd - Inches(0.8), Inches(0.5),
         [{'runs': [{'t': '✓   ', 'size': 15, 'color': GREEN, 'bold': True}, {'t': it, 'size': 15, 'color': WHITE}]}])
    yy += Inches(0.62)
notes(s, "LAYOUT: Two-column before/after compare. ANIMATION: left card Fade, right card Float-in then a subtle Pulse to draw the eye. Right card is highlighted with blue border.")

# ============================================================================= SLIDE 6
s = slide(); grad_bg(s, angle=30); footer(s, 6)
text(s, 0, Inches(2.2), SW, Inches(0.5),
     [{'t': 'THE SHIFT', 'size': 13, 'color': BLUE_BR, 'bold': True, 'align': CENTER}], align=CENTER, track=4)
text(s, Inches(1.5), Inches(2.9), SW - Inches(3.0), Inches(1.8),
     [{'runs': [{'t': 'From passive recording\nto ', 'size': 40, 'color': WHITE, 'bold': True, 'font': TITLE},
                {'t': 'real-time prevention.', 'size': 40, 'color': BLUE_BR, 'bold': True, 'font': TITLE}], 'align': CENTER, 'ls': 1.1}],
     align=CENTER, anchor=MIDDLE)
notes(s, "LAYOUT: Full-bleed section divider. ELEMENT: one bold statement, key phrase in blue. ANIMATION: Morph transition in; text Fade. Pause here — this is the pivot of the story.")

# ============================================================================= SLIDE 7
s = slide(); header(s, "The Solution", "One platform. Four ways it protects you.", 7)
pill = [("01", "Detect", "Pose + object AI flags suspicious behaviour live."),
        ("02", "Alert", "Instant siren, Telegram and email with a snapshot."),
        ("03", "Analyze", "Upload any video for an annotated forensic review."),
        ("04", "Recognize", "Blacklist & VIP faces matched automatically.")]
gap = Inches(0.32); cwd = (CW - 3 * gap) / 4
for i, (n, h, sub) in enumerate(pill):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.5), cwd, Inches(3.2))
    chip(s, x + Inches(0.3), Inches(2.85), Inches(0.7), n)
    text(s, x + Inches(0.3), Inches(3.85), cwd - Inches(0.6), Inches(0.5),
         [{'t': h, 'size': 19, 'color': WHITE, 'bold': True}])
    text(s, x + Inches(0.3), Inches(4.4), cwd - Inches(0.6), Inches(1.3),
         [{'t': sub, 'size': 13, 'color': MUTED, 'ls': 1.05}])
notes(s, "LAYOUT: 4 pillar cards. ELEMENT: numbered chips. ANIMATION: cards Fade-in staggered L→R. These four verbs structure the rest of the deck.")

# ============================================================================= SLIDE 8
s = slide(); grad_bg(s, angle=120); footer(s, 8)
text(s, Inches(1.3), Inches(1.2), Inches(2), Inches(2),
     [{'t': '“', 'size': 120, 'color': BLUE, 'bold': True, 'font': TITLE}])
text(s, Inches(1.6), Inches(2.6), SW - Inches(3.2), Inches(2.2),
     [{'t': 'Make every camera a proactive security guard — one that never blinks, never tires, and acts the instant it matters.',
       'size': 30, 'color': WHITE, 'bold': True, 'font': TITLE, 'ls': 1.15}])
text(s, Inches(1.6), Inches(5.2), Inches(8), Inches(0.4),
     [{'t': 'OUR VISION', 'size': 13, 'color': BLUE_BR, 'bold': True}], track=3)
notes(s, "LAYOUT: Quote / vision slide. ELEMENT: oversized quotation mark. ANIMATION: text Fade-up; keep minimal. Delivered as the mission statement.")

# ============================================================================= SLIDE 9
s = slide(); header(s, "How It Works", "Three steps, milliseconds apart", 9)
steps = [("Watch", "Multi-threaded engine ingests every camera feed in parallel."),
         ("Understand", "YOLOv8 pose & object models read posture, gestures and zones."),
         ("Alert", "Risk crosses threshold → siren + remote notification fire.")]
gap = Inches(0.55); cwd = (CW - 2 * gap) / 3
for i, (h, sub) in enumerate(steps):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.7), cwd, Inches(2.9))
    text(s, x + Inches(0.35), Inches(2.55), Inches(2), Inches(1),
         [{'t': f'0{i+1}', 'size': 46, 'color': BLUE, 'bold': True, 'font': TITLE}])
    text(s, x + Inches(0.35), Inches(3.7), cwd - Inches(0.7), Inches(0.5),
         [{'t': h, 'size': 21, 'color': WHITE, 'bold': True}])
    text(s, x + Inches(0.35), Inches(4.25), cwd - Inches(0.7), Inches(1.2),
         [{'t': sub, 'size': 13.5, 'color': MUTED, 'ls': 1.1}])
    if i < 2:
        arrow(s, x + cwd + Inches(0.07), Inches(3.85), Inches(0.4), Inches(0.55))
notes(s, "LAYOUT: 3-step horizontal stepper with chevrons. ANIMATION: reveal step 1→2→3 with Wipe (L→R); chevrons Fade between. Narrate the pipeline.")

# ============================================================================= SLIDE 10
s = slide(); header(s, "Architecture", "A clean pipeline from lens to alert", 10)
flow = [("CAMERAS", "Webcam / RTSP\nmulti-feed"), ("AI ENGINE", "YOLOv8 Pose\n+ Object"),
        ("DASHBOARD", "Live view\n+ analytics"), ("ALERTS", "Siren · Telegram\n· Email")]
gap = Inches(0.62); bw = (CW - 3 * gap) / 4
for i, (h, sub) in enumerate(flow):
    x = MARGIN + i * (bw + gap)
    fill = CARD_HI if i in (1,) else CARD
    line = BLUE if i == 1 else CARD_LINE
    card(s, x, Inches(2.9), bw, Inches(1.9), fill=fill, line=line, lw=1.5 if i == 1 else 1.0)
    text(s, x, Inches(3.15), bw, Inches(0.4), [{'t': h, 'size': 13, 'color': BLUE_BR, 'bold': True, 'align': CENTER}], align=CENTER, track=1.5)
    text(s, x, Inches(3.7), bw, Inches(0.9), [{'t': sub, 'size': 14, 'color': LIGHT, 'align': CENTER, 'ls': 1.05}], align=CENTER)
    if i < 3:
        arrow(s, x + bw + Inches(0.12), Inches(3.55), Inches(0.38), Inches(0.6))
text(s, MARGIN, Inches(5.4), CW, Inches(0.5),
     [{'t': 'FastAPI + WebSockets stream results in real time · SQLite stores every alert & face encoding.',
       'size': 13, 'color': MUTED, 'align': CENTER}], align=CENTER)
notes(s, "LAYOUT: Left→right flow diagram, AI engine emphasized. ANIMATION: boxes appear in sequence, chevrons Wipe between. Mention WebSocket + SQLite at the end.")

# ============================================================================= SLIDE 11
s = slide(); header(s, "Capability · Performance", "Multi-threaded camera engine", 11)
text(s, MARGIN, Inches(2.5), Inches(6.0), Inches(3.5),
     [{'t': 'Zero-lag, every feed at once', 'size': 20, 'color': WHITE, 'bold': True, 'sa': 10},
      {'t': 'Each camera runs on its own capture thread, so frames never queue behind one another.',
       'size': 15, 'color': LIGHT, 'ls': 1.25, 'sa': 8},
      {'t': 'Tracker states are isolated per (camera, person) pair — no identity bleed across feeds, even with dozens of people in frame.',
       'size': 15, 'color': LIGHT, 'ls': 1.25}])
# right: parallel thread bars feeding a core
bx = Inches(7.6)
for i in range(4):
    yy = Inches(2.55) + i * Inches(0.78)
    card(s, bx, yy, Inches(3.0), Inches(0.6), fill=CARD, line=CARD_LINE)
    text(s, bx + Inches(0.25), yy, Inches(2.6), Inches(0.6),
         [{'runs': [{'t': f'CAM 0{i+1}', 'size': 12, 'color': MUTED, 'bold': True}]}], anchor=MIDDLE)
    rect(s, bx + Inches(2.1), yy + Inches(0.21), Inches(0.7), Inches(0.18), fill=BLUE if i % 2 == 0 else BLUE_BR)
    arrow(s, bx + Inches(3.05), yy + Inches(0.12), Inches(0.35), Inches(0.36), color=CARD_LINE)
card(s, bx + Inches(3.5), Inches(3.35), Inches(1.6), Inches(1.3), fill=CARD_HI, line=BLUE, lw=1.5)
text(s, bx + Inches(3.5), Inches(3.55), Inches(1.6), Inches(1.0),
     [{'t': 'AI\nCORE', 'size': 16, 'color': WHITE, 'bold': True, 'align': CENTER, 'ls': 1.0}], align=CENTER, anchor=MIDDLE)
notes(s, "LAYOUT: text left, schematic right (parallel threads → shared core). ANIMATION: CAM bars Wipe-in top→bottom, AI CORE Zoom. Stresses concurrency + no identity bleed.")

# ============================================================================= SLIDE 12
s = slide(); header(s, "Capability · Intelligence", "It reads behaviour, not just pixels", 12)
feats = [("Concealment", "Tracks hand-to-pocket / bag gestures after picking up an item."),
         ("Loitering", "Flags dwell time beyond a configurable threshold in any zone."),
         ("Zone intrusion", "Sirens when wrists cross restricted ROI boundaries."),
         ("Suspicious posture", "Detects sudden bending in low-visibility aisles.")]
gap = Inches(0.4); cwd = (CW - gap) / 2; ch = Inches(1.75)
for i, (h, sub) in enumerate(feats):
    col, row = i % 2, i // 2
    x = MARGIN + col * (cwd + gap); y = Inches(2.5) + row * (ch + Inches(0.35))
    card(s, x, y, cwd, ch)
    chip(s, x + Inches(0.3), y + Inches(0.32), Inches(0.55), str(i + 1), fill=BLUE, fsize=14)
    text(s, x + Inches(1.05), y + Inches(0.28), cwd - Inches(1.4), Inches(0.5),
         [{'t': h, 'size': 18, 'color': WHITE, 'bold': True}])
    text(s, x + Inches(1.05), y + Inches(0.82), cwd - Inches(1.4), Inches(0.8),
         [{'t': sub, 'size': 13, 'color': MUTED, 'ls': 1.1}])
notes(s, "LAYOUT: 2x2 feature grid. ANIMATION: cards Fade-in (Z pattern). Each is a distinct behavioural model — keep it skimmable.")

# ============================================================================= SLIDE 13
s = slide(); header(s, "Capability · Deep Dive", "Concealment & loitering, decoded", 13)
gap = Inches(0.5); cwd = (CW - gap) / 2
# concealment
card(s, MARGIN, Inches(2.5), cwd, Inches(3.5))
text(s, MARGIN + Inches(0.4), Inches(2.8), cwd - Inches(0.8), Inches(0.5),
     [{'t': 'CONCEALMENT', 'size': 13, 'color': BLUE_BR, 'bold': True}], track=2)
text(s, MARGIN + Inches(0.4), Inches(3.3), cwd - Inches(0.8), Inches(2.4),
     [{'t': 'Item picked up', 'size': 17, 'color': WHITE, 'bold': True, 'sa': 4},
      {'t': '+ hand moves to pocket / bag', 'size': 15, 'color': LIGHT, 'sa': 4},
      {'t': '+ item no longer visible', 'size': 15, 'color': LIGHT, 'sa': 12},
      {'t': '= concealment flagged', 'size': 17, 'color': GREEN, 'bold': True}])
# loitering
x2 = MARGIN + cwd + gap
card(s, x2, Inches(2.5), cwd, Inches(3.5))
text(s, x2 + Inches(0.4), Inches(2.8), cwd - Inches(0.8), Inches(0.5),
     [{'t': 'LOITERING', 'size': 13, 'color': BLUE_BR, 'bold': True}], track=2)
text(s, x2 + Inches(0.4), Inches(3.3), cwd - Inches(0.8), Inches(0.5),
     [{'t': 'Dwell time vs. threshold', 'size': 17, 'color': WHITE, 'bold': True}])
# progress bar
rect(s, x2 + Inches(0.4), Inches(4.1), cwd - Inches(0.8), Inches(0.32), fill=CARD_LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
rect(s, x2 + Inches(0.4), Inches(4.1), (cwd - Inches(0.8)) * 0.72, Inches(0.32), fill=BLUE_BR, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
text(s, x2 + Inches(0.4), Inches(4.6), cwd - Inches(0.8), Inches(1.2),
     [{'t': 'Cross the line and the alarm trips automatically — tune the seconds per zone from the dashboard.',
       'size': 13.5, 'color': MUTED, 'ls': 1.15}])
notes(s, "LAYOUT: two deep-dive cards. ANIMATION: concealment lines reveal one-by-one (Appear) to show the logic chain; loitering bar Wipe L→R to threshold.")

# ============================================================================= SLIDE 14
s = slide(); header(s, "Capability · Zones", "Draw the line. We watch it.", 14)
text(s, MARGIN, Inches(2.5), Inches(5.7), Inches(3.4),
     [{'t': 'Interactive ROI drawer', 'size': 20, 'color': WHITE, 'bold': True, 'sa': 10},
      {'t': 'Draw security polygons straight onto the live feed inside the dashboard.',
       'size': 15, 'color': LIGHT, 'ls': 1.25, 'sa': 8},
      {'t': 'Coordinates auto-scale to the 1280×720 matrix, so zones stay accurate on any screen — and the instant a wrist crosses in, the siren fires.',
       'size': 15, 'color': LIGHT, 'ls': 1.25}])
# stylised feed with ROI polygon
fx, fy, fw, fh = Inches(7.1), Inches(2.45), Inches(5.3), Inches(3.35)
card(s, fx, fy, fw, fh, fill=CARD_HI, line=CARD_LINE)
poly = s.shapes.add_shape(MSO_SHAPE.HEXAGON, fx + Inches(1.1), fy + Inches(0.7), Inches(3.0), Inches(2.0))
poly.fill.solid(); poly.fill.fore_color.rgb = BLUE
try:
    poly.fill.transparency = 0  # base
except Exception:
    pass
poly.line.color.rgb = BLUE_BR; poly.line.width = Pt(2); _noshadow(poly)
text(s, fx + Inches(1.1), fy + Inches(1.4), Inches(3.0), Inches(0.5),
     [{'t': 'RESTRICTED ZONE', 'size': 12, 'color': WHITE, 'bold': True, 'align': CENTER}], align=CENTER)
text(s, fx, fy + fh - Inches(0.45), fw, Inches(0.4),
     [{'t': 'live feed · 1280 × 720', 'size': 10, 'color': MUTED, 'align': CENTER}], align=CENTER)
notes(s, "LAYOUT: text left, stylised live-feed-with-ROI right. ANIMATION: polygon draws on (Wipe) then a red Pulse to simulate an intrusion trip. Real ROI UI shown on slide 19.")

# ============================================================================= SLIDE 15
s = slide(); header(s, "Capability · Identity", "Faces that matter, matched instantly", 15)
gap = Inches(0.5); cwd = (CW - gap) / 2
# blacklist
card(s, MARGIN, Inches(2.5), cwd, Inches(3.4), fill=CARD, line=RED, lw=1.3)
chip(s, MARGIN + Inches(0.4), Inches(2.85), Inches(0.7), '!', fill=RED, fsize=22)
text(s, MARGIN + Inches(0.4), Inches(3.8), cwd - Inches(0.8), Inches(0.5),
     [{'t': 'Blacklist', 'size': 22, 'color': WHITE, 'bold': True}])
text(s, MARGIN + Inches(0.4), Inches(4.35), cwd - Inches(0.8), Inches(1.3),
     [{'t': 'Known offenders trigger a high-priority alarm and auto-record evidence the moment they enter.',
       'size': 14.5, 'color': LIGHT, 'ls': 1.2}])
# vip
x2 = MARGIN + cwd + gap
card(s, x2, Inches(2.5), cwd, Inches(3.4), fill=CARD, line=GREEN, lw=1.3)
chip(s, x2 + Inches(0.4), Inches(2.85), Inches(0.7), '★', fill=GREEN, fsize=20)
text(s, x2 + Inches(0.4), Inches(3.8), cwd - Inches(0.8), Inches(0.5),
     [{'t': 'VIP Whitelist', 'size': 22, 'color': WHITE, 'bold': True}])
text(s, x2 + Inches(0.4), Inches(4.35), cwd - Inches(0.8), Inches(1.3),
     [{'t': 'Trusted staff and loyal customers are recognised and welcomed with a green badge — no false alarms.',
       'size': 14.5, 'color': LIGHT, 'ls': 1.2}])
notes(s, "LAYOUT: split identity cards (red vs green). ANIMATION: both Float-in; blacklist card subtle red Pulse. Built on Dlib face encodings stored in SQLite.")

# ============================================================================= SLIDE 16
s = slide(); header(s, "Capability · Response", "The alarm reaches you anywhere", 16)
chans = [("Browser siren", "A synthesised emergency siren sweeps instantly — no audio files, no lag."),
         ("Telegram", "Bot pushes the alert + snapshot to your phone in seconds."),
         ("Email (SMTP)", "Automated mail with evidence for the record and the team.")]
gap = Inches(0.4); cwd = (CW - 2 * gap) / 3
for i, (h, sub) in enumerate(chans):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.6), cwd, Inches(3.0))
    chip(s, x + Inches(0.35), Inches(2.95), Inches(0.65), ['♪', '✈', '@'][i], fill=BLUE, fsize=18)
    text(s, x + Inches(0.35), Inches(3.85), cwd - Inches(0.7), Inches(0.5),
         [{'t': h, 'size': 18, 'color': WHITE, 'bold': True}])
    text(s, x + Inches(0.35), Inches(4.4), cwd - Inches(0.7), Inches(1.2),
         [{'t': sub, 'size': 13.5, 'color': MUTED, 'ls': 1.15}])
text(s, MARGIN, Inches(6.0), CW, Inches(0.4),
     [{'t': 'Smart 3-second cooldown prevents alert fatigue.', 'size': 13, 'color': BLUE_BR, 'align': CENTER}], align=CENTER)
notes(s, "LAYOUT: 3-channel triptych. ANIMATION: cards Fade-in L→R, cooldown line Fade last. Emphasise multi-channel + anti-fatigue cooldown.")

# ============================================================================= SLIDE 17
s = slide(); header(s, "Capability · Forensics", "Upload a video. Get answers.", 17)
text(s, MARGIN, Inches(2.5), Inches(3.9), Inches(3.6),
     [{'t': 'Any MP4, fully analysed', 'size': 19, 'color': WHITE, 'bold': True, 'sa': 10},
      {'t': 'Drop in recorded footage and the same AI runs frame-by-frame.', 'size': 14, 'color': LIGHT, 'ls': 1.2, 'sa': 8},
      {'t': '· Annotated video returned\n· Every alert listed\n· Telegram notification on detection',
       'size': 14, 'color': LIGHT, 'ls': 1.5}])
shot(s, f"{SHOTS}/video-analysis.png", Inches(5.1), Inches(1.95), Inches(7.4))
notes(s, "LAYOUT: text left, product screenshot right (browser frame). ANIMATION: screenshot Float-in (Up); bullet list Fade-up. This is a standout differentiator — demo it live if possible.")

# ============================================================================= SLIDE 18
s = slide(); solid_bg(s); top_bar(s); footer(s, 18)
text(s, MARGIN, Inches(0.62), Inches(11), Inches(0.35),
     [{'t': 'THE PRODUCT', 'size': 12, 'color': BLUE_BR, 'bold': True}], track=2.5)
text(s, MARGIN, Inches(1.0), Inches(11.5), Inches(0.7),
     [{'t': 'One glassmorphic control room', 'size': 30, 'color': WHITE, 'bold': True, 'font': TITLE}])
shot(s, f"{SHOTS}/overview.png", Inches(2.46), Inches(1.75), Inches(8.4))
notes(s, "LAYOUT: full-bleed hero screenshot, minimal title. ANIMATION: screenshot Float-in / Zoom slightly. Let the product breathe — least text, most impact.")

# ============================================================================= SLIDES 19-22 (screenshot + caption)
shot_slides = [
    (19, "Camera Setup & ROI", "cameras.png", "Add webcams or RTSP streams and draw security zones — saved per camera in cameras.json."),
    (20, "Face Management", "faces.png", "Upload portraits, tag Blacklist or VIP, and delete instantly with live DB sync."),
    (21, "Alerts & Analytics", "history.png", "Search, filter by category, watch weekly trends and export the log to CSV in one click."),
    (22, "Notifications", "settings.png", "Wire up Telegram and SMTP in minutes — test and broadcast alerts with snapshots."),
]
for n, title, img, cap in shot_slides:
    s = slide(); header(s, "The Product", title, n, tsize=30)
    text(s, MARGIN, Inches(2.4), Inches(3.7), Inches(3.0),
         [{'t': cap, 'size': 15, 'color': LIGHT, 'ls': 1.3}])
    shot(s, f"{SHOTS}/{img}", Inches(5.0), Inches(2.05), Inches(7.5))
    notes(s, f"LAYOUT: caption left, screenshot right in browser frame. ANIMATION: screenshot Float-in; caption Fade-up. Real UI of '{title}'.")

# ============================================================================= SLIDE 23
s = slide(); header(s, "Under the Hood", "Built on a proven, modern stack", 23)
groups = [("AI / VISION", ["YOLOv8 Pose", "YOLOv8 Object", "OpenCV", "Dlib face encodings"]),
          ("BACKEND", ["Python 3.10+", "FastAPI", "WebSockets", "SQLite"]),
          ("FRONTEND", ["Next.js 16", "React 19", "Tailwind CSS", "Recharts"])]
gap = Inches(0.4); cwd = (CW - 2 * gap) / 3
for i, (h, items) in enumerate(groups):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.5), cwd, Inches(3.3))
    text(s, x + Inches(0.4), Inches(2.8), cwd - Inches(0.8), Inches(0.4),
         [{'t': h, 'size': 13, 'color': BLUE_BR, 'bold': True}], track=2)
    yy = Inches(3.45)
    for it in items:
        rect(s, x + Inches(0.4), yy + Inches(0.09), Inches(0.12), Inches(0.12), fill=BLUE, shape=MSO_SHAPE.OVAL)
        text(s, x + Inches(0.7), yy, cwd - Inches(1.1), Inches(0.45),
             [{'t': it, 'size': 15, 'color': LIGHT}])
        yy += Inches(0.55)
notes(s, "LAYOUT: 3 grouped tech columns. ANIMATION: columns Fade-in L→R, list items Appear. Signals engineering maturity without jargon overload.")

# ============================================================================= SLIDE 24
s = slide(); grad_bg(s, angle=45); footer(s, 24)
text(s, 0, Inches(2.0), SW, Inches(0.5),
     [{'t': 'SEE IT LIVE', 'size': 14, 'color': BLUE_BR, 'bold': True, 'align': CENTER}], align=CENTER, track=4)
text(s, 0, Inches(2.7), SW, Inches(1.0),
     [{'t': 'Try the dashboard now', 'size': 40, 'color': WHITE, 'bold': True, 'font': TITLE, 'align': CENTER}], align=CENTER)
card(s, Inches(3.4), Inches(4.2), Inches(6.5), Inches(0.95), fill=CARD_HI, line=BLUE, lw=1.5, radius=0.5)
text(s, Inches(3.4), Inches(4.2), Inches(6.5), Inches(0.95),
     [{'t': 'theft-detection-dusky.vercel.app', 'size': 22, 'color': WHITE, 'bold': True, 'align': CENTER}],
     align=CENTER, anchor=MIDDLE)
notes(s, "LAYOUT: divider CTA. ELEMENT: pill-shaped URL chip (add a QR code in PowerPoint for live audiences). ANIMATION: Morph in; URL chip Zoom. Invite the room to open it on their phones.")

# ============================================================================= SLIDE 25
s = slide(); header(s, "Differentiation", "Why we win", 25)
rows = [["", "Legacy CCTV", "Cloud SaaS", "Theft Detection"],
        ["Real-time prevention", "✕", "~", "✓"],
        ["Behaviour + pose AI", "✕", "~", "✓"],
        ["Runs on-premise", "✓", "✕", "✓"],
        ["Video forensic upload", "✕", "~", "✓"],
        ["One-time, no per-seat fee", "✓", "✕", "✓"]]
tx, ty = MARGIN, Inches(2.4)
tw = CW; rh = Inches(0.62); c0 = Inches(4.6)
cw = (tw - c0) / 3
for r, row in enumerate(rows):
    y = ty + r * rh
    head = (r == 0)
    if head:
        rect(s, tx, y, tw, rh, fill=CARD_HI)
    elif r % 2 == 1:
        rect(s, tx, y, tw, rh, fill=CARD)
    # col 0
    text(s, tx + Inches(0.3), y, c0 - Inches(0.4), rh,
         [{'t': row[0], 'size': 14, 'color': WHITE if head else LIGHT, 'bold': head}], anchor=MIDDLE)
    for c in range(1, 4):
        cx = tx + c0 + (c - 1) * cw
        val = row[c]
        col = WHITE
        if val == '✓': col = GREEN
        elif val == '✕': col = RED
        elif val == '~': col = AMBER
        hi = (c == 3 and not head)
        if c == 3:
            rect(s, cx, y, cw, rh, fill=None, line=BLUE, lw=1.0) if head else rect(s, cx, y, cw, rh, fill=CARD_HI)
        text(s, cx, y, cw, rh,
             [{'t': val, 'size': 16 if not head else 13, 'color': BLUE_BR if head and c == 3 else (col if not head else MUTED),
               'bold': True, 'align': CENTER}], align=CENTER, anchor=MIDDLE)
rect(s, tx, ty, tw, rh * len(rows), fill=None, line=CARD_LINE, lw=1.0)
notes(s, "LAYOUT: comparison matrix, our column highlighted. ANIMATION: reveal row-by-row (Wipe down); our column Fade emphasis. ✓ full · ~ partial · ✕ none.")

# ============================================================================= SLIDE 26
s = slide(); header(s, "Opportunity", "A large, growing market", 26)
text(s, MARGIN, Inches(2.4), Inches(5.5), Inches(3),
     [{'t': '$40B+', 'size': 72, 'color': BLUE_BR, 'bold': True, 'font': TITLE, 'sa': 4},
      {'t': 'global video surveillance market', 'size': 16, 'color': LIGHT, 'sa': 14},
      {'t': 'Driven by AI adoption, retail shrink pressure and the shift from passive to smart cameras.',
       'size': 14, 'color': MUTED, 'ls': 1.25}])
# funnel TAM/SAM/SOM
fx = Inches(7.4); widths = [Inches(5.1), Inches(3.9), Inches(2.7)]
labels = [("TAM", "All video surveillance"), ("SAM", "AI retail loss-prevention"), ("SOM", "SMB stores · year 1-3")]
yy = Inches(2.5)
for i, (lab, sub) in enumerate(labels):
    w = widths[i]; x = fx + (Inches(5.1) - w) / 2
    card(s, x, yy, w, Inches(0.95), fill=CARD_HI if i == 2 else CARD, line=BLUE if i == 2 else CARD_LINE, lw=1.4 if i == 2 else 1.0)
    text(s, x, yy, w, Inches(0.95),
         [{'runs': [{'t': lab + '  ', 'size': 16, 'color': BLUE_BR, 'bold': True}, {'t': sub, 'size': 12.5, 'color': LIGHT}], 'align': CENTER}],
         align=CENTER, anchor=MIDDLE)
    yy += Inches(1.15)
text(s, MARGIN, Inches(6.4), CW, Inches(0.3),
     [{'t': '*Illustrative — insert your validated market sizing.', 'size': 10, 'color': MUTED}])
notes(s, "LAYOUT: big stat left, TAM/SAM/SOM funnel right. ANIMATION: stat Zoom; funnel tiers Wipe top→bottom narrowing focus. *replace figures with sourced data.")

# ============================================================================= SLIDE 27
s = slide(); header(s, "Business Model", "Simple, scalable pricing", 27)
tiers = [("Starter", "Single store", ["1 location", "Up to 4 cameras", "Core detection", "Email alerts"], False),
         ("Pro", "Growing retail", ["Up to 5 locations", "Unlimited cameras", "Face recognition", "Telegram + analytics"], True),
         ("Enterprise", "Chains & malls", ["Unlimited locations", "On-prem / private cloud", "Custom AI model", "Priority support"], False)]
gap = Inches(0.4); cwd = (CW - 2 * gap) / 3
for i, (name, who, feats, hi) in enumerate(tiers):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.35), cwd, Inches(3.85), fill=CARD_HI if hi else CARD, line=BLUE if hi else CARD_LINE, lw=1.8 if hi else 1.0)
    if hi:
        rect(s, x + cwd / 2 - Inches(0.8), Inches(2.18), Inches(1.6), Inches(0.34), fill=BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        text(s, x + cwd / 2 - Inches(0.8), Inches(2.18), Inches(1.6), Inches(0.34),
             [{'t': 'POPULAR', 'size': 10, 'color': WHITE, 'bold': True, 'align': CENTER}], align=CENTER, anchor=MIDDLE, track=1.5)
    text(s, x + Inches(0.35), Inches(2.65), cwd - Inches(0.7), Inches(0.5),
         [{'t': name, 'size': 22, 'color': WHITE, 'bold': True}])
    text(s, x + Inches(0.35), Inches(3.15), cwd - Inches(0.7), Inches(0.4),
         [{'t': who, 'size': 13, 'color': BLUE_BR}])
    yy = Inches(3.75)
    for f in feats:
        text(s, x + Inches(0.35), yy, cwd - Inches(0.7), Inches(0.45),
             [{'runs': [{'t': '✓  ', 'size': 13, 'color': GREEN, 'bold': True}, {'t': f, 'size': 13.5, 'color': LIGHT}]}])
        yy += Inches(0.5)
notes(s, "LAYOUT: 3 pricing tiers, middle highlighted with POPULAR ribbon. ANIMATION: cards Float-in, middle scales up slightly. Tiers map to the Business Model Canvas in the report.")

# ============================================================================= SLIDE 28
s = slide(); header(s, "Financials", "The numbers behind the build", 28)
kpis = [("980K", "Assets (DZD)", BLUE_BR), ("6.0M", "Revenue (DZD)", GREEN),
        ("3.3M", "Total cost (DZD)", AMBER), ("2.7M", "Net cash flow (DZD)", WHITE)]
gap = Inches(0.35); cwd = (CW - 3 * gap) / 4
for i, (num, lab, col) in enumerate(kpis):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.7), cwd, Inches(2.4))
    text(s, x, Inches(3.1), cwd, Inches(1.0),
         [{'t': num, 'size': 40, 'color': col, 'bold': True, 'font': TITLE, 'align': CENTER}], align=CENTER)
    text(s, x + Inches(0.2), Inches(4.25), cwd - Inches(0.4), Inches(0.7),
         [{'t': lab, 'size': 13, 'color': MUTED, 'align': CENTER}], align=CENTER)
text(s, MARGIN, Inches(5.6), CW, Inches(0.4),
     [{'t': 'Projected first-cycle figures — full tables in the Final Year Report (Ch. 3).', 'size': 12, 'color': MUTED, 'align': CENTER}], align=CENTER)
notes(s, "LAYOUT: 4 KPI cards. ANIMATION: numbers Zoom/count-up staggered. Tie back to the detailed financial chapter; net positive is the headline.")

# ============================================================================= SLIDE 29
s = slide(); header(s, "Roadmap", "Where we go next", 29)
phases = [("NOW", "Live & shipping", ["Real-time detection", "Dashboard + alerts", "Video analysis"]),
          ("NEXT", "In progress", ["Custom-trained model", "Multi-store cloud sync", "Role-based access"]),
          ("LATER", "Vision", ["Mobile companion app", "Predictive heatmaps", "POS integration"])]
# timeline line
rect(s, MARGIN + Inches(0.2), Inches(3.05), CW - Inches(0.4), Pt(2.5), fill=CARD_LINE)
gap = Inches(0.4); cwd = (CW - 2 * gap) / 3
for i, (ph, sub, items) in enumerate(phases):
    x = MARGIN + i * (cwd + gap)
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), Inches(2.92), Inches(0.28), Inches(0.28))
    d.fill.solid(); d.fill.fore_color.rgb = BLUE_BR if i == 0 else BLUE; d.line.color.rgb = NAVY; d.line.width = Pt(3); _noshadow(d)
    card(s, x, Inches(3.5), cwd, Inches(2.7))
    text(s, x + Inches(0.35), Inches(3.75), cwd - Inches(0.7), Inches(0.4),
         [{'t': ph, 'size': 14, 'color': BLUE_BR, 'bold': True}], track=2)
    text(s, x + Inches(0.35), Inches(4.15), cwd - Inches(0.7), Inches(0.4),
         [{'t': sub, 'size': 13, 'color': MUTED}])
    yy = Inches(4.7)
    for it in items:
        text(s, x + Inches(0.35), yy, cwd - Inches(0.7), Inches(0.4),
             [{'runs': [{'t': '· ', 'size': 14, 'color': BLUE, 'bold': True}, {'t': it, 'size': 13.5, 'color': LIGHT}]}])
        yy += Inches(0.45)
notes(s, "LAYOUT: 3-phase horizontal timeline (Now/Next/Later) on a track. ANIMATION: nodes light up L→R (Appear), cards Float-in. 'Now' node is brightest to anchor credibility.")

# ============================================================================= SLIDE 30
s = slide(); grad_bg(s, angle=60); footer(s, 30)
radar(s, Inches(11.6), Inches(5.6))
try:
    s.shapes.add_picture(LOGO, MARGIN, Inches(0.8), height=Inches(0.95))
except Exception:
    pass
text(s, MARGIN, Inches(2.4), Inches(10), Inches(0.5),
     [{'t': 'THANK YOU', 'size': 14, 'color': BLUE_BR, 'bold': True}], track=4)
text(s, MARGIN, Inches(2.95), Inches(10.5), Inches(1.6),
     [{'t': 'Let’s make every camera\nsee what matters.', 'size': 40, 'color': WHITE, 'bold': True, 'font': TITLE, 'ls': 1.1}])
rect(s, MARGIN, Inches(4.9), Inches(0.8), Pt(3), fill=BLUE)
text(s, MARGIN, Inches(5.4), Inches(11), Inches(1.2),
     [{'runs': [{'t': 'Live demo   ', 'size': 14, 'color': MUTED}, {'t': 'theft-detection-dusky.vercel.app', 'size': 14, 'color': BLUE_BR, 'bold': True}]}])
notes(s, "LAYOUT: closing hero, mirrors the title. ELEMENT: logo, demo link, radar motif. ANIMATION: Morph from previous; text Fade-up. End on the one-line vision and the live URL.")

# ----------------------------------------------------------------------------- save + verify
out = "Theft_Detection_Pitch.pptx"
prs.save(out)
chk = Presentation(out)
print(f"Saved {out}  |  slides: {len(chk.slides.__iter__.__self__._sldIdLst)}  |  size: {os.path.getsize(out)//1024} KB")
