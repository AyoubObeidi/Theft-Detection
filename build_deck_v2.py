# -*- coding: utf-8 -*-
"""
Theft Detection — pitch deck v2.
Distinct "Design System B" (navy + electric cyan), denser content,
and REAL baked-in PowerPoint animations (Morph/Fade transitions +
auto-playing staggered fade-in entrance builds, injected as OOXML).

Run:  py build_deck_v2.py   ->  Theft_Detection_Pitch_v2.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# --------------------------------------------------------------------- palette B
BG     = RGBColor(0x0B, 0x15, 0x26)
BG2    = RGBColor(0x0E, 0x1B, 0x30)
GRAD2  = RGBColor(0x10, 0x2A, 0x4A)
CARD   = RGBColor(0x13, 0x24, 0x3F)
CARD2  = RGBColor(0x16, 0x2C, 0x4D)
LINE   = RGBColor(0x27, 0x3A, 0x5C)
BLUE   = RGBColor(0x2E, 0x74, 0xB5)
CYAN   = RGBColor(0x22, 0xD3, 0xEE)
CYAN_D = RGBColor(0x0E, 0x74, 0x90)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT   = RGBColor(0xE6, 0xED, 0xF6)
MUTE   = RGBColor(0x93, 0xA6, 0xC0)
GHOST  = RGBColor(0x12, 0x22, 0x3B)
GREEN  = RGBColor(0x34, 0xD3, 0x99)
RED    = RGBColor(0xFF, 0x6B, 0x6B)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)

BODY  = "Segoe UI"
TITLE = "Segoe UI Semibold"
LEFT, CENTER, RIGHT = PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT
TOP, MIDDLE = MSO_ANCHOR.TOP, MSO_ANCHOR.MIDDLE

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.95)
CW = SW - 2 * MARGIN
SHOTS = "docs/shots"
LOGO = "docs/djelfa_logo.png"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH

_buf = []      # shapes created since last mark (candidates for one animation step)
_steps = []    # list of step groups (each = list of shapes)
CURN = [0]     # running slide number (auto-assigned in new_slide so deletions renumber)
TOTAL = 28     # total slides in the deck (shown in footer)


# --------------------------------------------------------------------- animation OOXML
class _Ctr:
    def __init__(self): self.n = 2
    def nx(self):
        self.n += 1
        return self.n


def _effect_par(spid, delay, node, c):
    a, b, cc, d, e = c.nx(), c.nx(), c.nx(), c.nx(), c.nx()
    return (
        f'<p:par><p:cTn id="{a}" fill="hold"><p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
        f'<p:childTnLst><p:par><p:cTn id="{b}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst><p:par><p:cTn id="{cc}" presetID="10" presetClass="entr" presetSubtype="0" '
        f'fill="hold" grpId="0" nodeType="{node}"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
        f'<p:set><p:cBhvr><p:cTn id="{d}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName>'
        f'</p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        f'<p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{e}" dur="160"/>'
        f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>'
        f'</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
    )


def add_transition(slide, kind="fade"):
    if kind == "morph":
        xml = (
            '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
            f'xmlns:p="{P_NS}" xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main">'
            '<mc:Choice Requires="p14"><p:transition spd="fast" p14:dur="350"><p14:morph option="byObject"/>'
            '</p:transition></mc:Choice><mc:Fallback><p:transition spd="fast"><p:fade/></p:transition>'
            '</mc:Fallback></mc:AlternateContent>'
        )
    elif kind == "push":
        xml = f'<p:transition xmlns:p="{P_NS}" spd="med"><p:push dir="l"/></p:transition>'
    else:
        xml = f'<p:transition xmlns:p="{P_NS}" spd="fast"><p:fade/></p:transition>'
    slide._element.append(parse_xml(xml))


def build_timing(slide, steps):
    groups = [[sh.shape_id for sh in g] for g in steps if g]
    if not groups:
        return
    c = _Ctr()
    pars = []
    for i, g in enumerate(groups):
        for j, spid in enumerate(g):
            if i == 0 and j == 0:
                node, delay = "withEffect", 60
            elif j == 0:
                node, delay = "afterEffect", 40
            else:
                node, delay = "withEffect", 0
            pars.append(_effect_par(spid, delay, node, c))
    cond = ('<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            '<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>')
    xml = (
        f'<p:timing xmlns:p="{P_NS}"><p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" '
        f'nodeType="tmRoot"><p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" '
        f'nodeType="mainSeq"><p:childTnLst>{"".join(pars)}</p:childTnLst>{cond}</p:cTn>{cond}</p:seq>'
        f'</p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
    )
    slide._element.append(parse_xml(xml))


# --------------------------------------------------------------------- core helpers
def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _buf.clear(); _steps.clear()
    CURN[0] += 1
    return s


def mark():
    if _buf:
        _steps.append(list(_buf)); _buf.clear()


def anim_reset():
    _buf.clear()


def finish(slide, transition="fade"):
    mark()
    add_transition(slide, transition)
    build_timing(slide, _steps)


def _reg(sh):
    _buf.append(sh)
    return sh


def _noshadow(sh):
    sh.shadow.inherit = False
    return sh


def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None, reg=True):
    sh = s.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    _noshadow(sh)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sh.adjustments[0] = radius
        except Exception: pass
    return _reg(sh) if reg else sh


def oval(s, l, t, w, h, fill, line=None, reg=True):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(1)
    _noshadow(sh)
    return _reg(sh) if reg else sh


def text(s, l, t, w, h, lines, align=LEFT, anchor=TOP, track=None, reg=True):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get('align', align)
        p._p.get_or_add_pPr().set('rtl', '0')
        if ln.get('sb') is not None: p.space_before = Pt(ln['sb'])
        if ln.get('sa') is not None: p.space_after = Pt(ln['sa'])
        if ln.get('ls') is not None: p.line_spacing = ln['ls']
        runs = ln['runs'] if 'runs' in ln else [ln]
        for rd in runs:
            r = p.add_run()
            r.text = '‪' + rd['t'] + '‬'      # LTR embed (locale bidi fix)
            r.font.size = Pt(rd.get('size', 14))
            r.font.bold = rd.get('bold', False)
            r.font.italic = rd.get('italic', False)
            r.font.name = rd.get('font', BODY)
            r.font.color.rgb = rd.get('color', TEXT)
            if rd.get('track'):
                try: r.font._rPr.set('spc', str(int(rd['track'] * 100)))
                except Exception: pass
        if track:
            for r in p.runs:
                try: r.font._rPr.set('spc', str(int(track * 100)))
                except Exception: pass
    return _reg(tb) if reg else tb


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def card(s, l, t, w, h, fill=CARD, line=LINE, lw=1.0, radius=0.045):
    return rect(s, l, t, w, h, fill=fill, line=line, lw=lw, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)


def chip(s, l, t, sz, label, fill=BLUE, fg=WHITE, fsize=15, radius=0.25):
    rect(s, l, t, sz, sz, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)
    text(s, l, t, sz, sz, [{'t': label, 'size': fsize, 'color': fg, 'bold': True, 'align': CENTER}],
         align=CENTER, anchor=MIDDLE)


def arrow(s, l, t, w, h, color=CYAN_D):
    a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); _noshadow(a)
    return _reg(a)


def metric(s, l, t, w, val, lab):
    card(s, l, t, w, Inches(0.95), fill=BG2, line=CYAN_D, lw=1.1)
    text(s, l + Inches(0.22), t + Inches(0.12), w - Inches(0.4), Inches(0.45),
         [{'t': val, 'size': 20, 'color': CYAN, 'bold': True, 'font': TITLE}])
    text(s, l + Inches(0.22), t + Inches(0.55), w - Inches(0.4), Inches(0.35),
         [{'t': lab, 'size': 10.5, 'color': MUTE}])


def shot(s, img, l, t, w, ar=1.6):
    bar = Inches(0.28); iw = w - Inches(0.08); ih = iw / ar
    card(s, l, t, w, bar + ih + Inches(0.06), fill=CARD2, line=LINE, radius=0.025)
    for i, c in enumerate([RED, AMBER, GREEN]):
        oval(s, l + Inches(0.20) + i * Inches(0.24), t + Inches(0.09), Inches(0.11), Inches(0.11), c)
    s.shapes.add_picture(img, l + Inches(0.04), t + bar, width=iw)  # picture registered via parent step
    _reg(s.shapes[-1])
    return t + bar + ih


# --------------------------------------------------------------------- chrome (static, not animated)
def bg(s, color=BG):
    rect(s, 0, 0, SW, SH, fill=color, reg=False)


def grad(s, c1=BG, c2=GRAD2, angle=60):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    sh.line.fill.background(); _noshadow(sh)
    try:
        sh.fill.gradient()
        st = sh.fill.gradient_stops
        st[0].position = 0.0; st[0].color.rgb = c1
        st[1].position = 1.0; st[1].color.rgb = c2
        try: sh.fill.gradient_angle = angle
        except Exception: pass
    except Exception:
        sh.fill.solid(); sh.fill.fore_color.rgb = c1


def dotgrid(s, x, y, cols, rows, gap=Inches(0.32), color=LINE):
    for r in range(rows):
        for c in range(cols):
            oval(s, x + c * gap, y + r * gap, Pt(2.2), Pt(2.2), color, reg=False)


def corner_brackets(s, color=CYAN_D):
    L = Inches(0.5); th = Pt(2)
    # top-left
    rect(s, MARGIN, Inches(0.55), L, th, fill=color, reg=False)
    rect(s, MARGIN, Inches(0.55), th, L, fill=color, reg=False)
    # bottom-right
    rect(s, SW - MARGIN - L, SH - Inches(0.57), L, th, fill=color, reg=False)
    rect(s, SW - MARGIN - th, SH - Inches(0.57) - L, th, L, fill=color, reg=False)


def ghostnum(s, n):
    text(s, SW - Inches(3.3), Inches(0.05), Inches(3.0), Inches(1.8),
         [{'t': f'{CURN[0]:02d}', 'size': 120, 'color': GHOST, 'bold': True, 'font': TITLE, 'align': RIGHT}],
         align=RIGHT, reg=False)


def chrome(s, n):
    rect(s, MARGIN, Inches(0.62), Inches(0.14), Inches(0.14), fill=CYAN, reg=False)
    text(s, MARGIN + Inches(0.26), Inches(0.58), Inches(5), Inches(0.3),
         [{'t': 'THEFT DETECTION', 'size': 9, 'color': MUTE, 'bold': True, 'track': 2}], reg=False)
    rect(s, MARGIN, SH - Inches(0.5), CW, Pt(0.75), fill=LINE, reg=False)
    text(s, MARGIN, SH - Inches(0.46), Inches(2), Inches(0.3),
         [{'t': f'{CURN[0]:02d} / {TOTAL}', 'size': 8.5, 'color': MUTE, 'bold': True, 'track': 1}], reg=False)
    text(s, SW - MARGIN - Inches(3), SH - Inches(0.46), Inches(3), Inches(0.3),
         [{'t': 'theft-detection-dusky.vercel.app', 'size': 8.5, 'color': MUTE, 'align': RIGHT}], align=RIGHT, reg=False)


def header(s, kicker, title, n, tsize=30):
    bg(s); chrome(s, n); ghostnum(s, n)
    rect(s, MARGIN, Inches(1.02), Inches(0.4), Pt(3), fill=CYAN, reg=False)
    text(s, MARGIN + Inches(0.55), Inches(0.92), Inches(10), Inches(0.35),
         [{'t': kicker.upper(), 'size': 11.5, 'color': CYAN, 'bold': True}], track=2.5)
    text(s, MARGIN, Inches(1.32), Inches(11.6), Inches(0.95),
         [{'t': title, 'size': tsize, 'color': WHITE, 'bold': True, 'font': TITLE}])
    anim_reset()   # discard chrome/title-rule from anim; we re-add title to its own step below
    # title becomes first animated step
    # (re-grab the title textbox: it's the last shape; register it)
    _reg(s.shapes[-1])
    mark()


def bullets(s, l, t, w, items, lead=13.5, sub=11.5, gap=6, marker=True):
    """items: list of (lead, subtext or None). Returns one textbox (one anim step)."""
    lines = []
    for i, (ld, sb) in enumerate(items):
        runs = []
        if marker:
            runs.append({'t': '› ', 'size': lead, 'color': CYAN, 'bold': True})
        runs.append({'t': ld, 'size': lead, 'color': TEXT, 'bold': True})
        lines.append({'runs': runs, 'sb': 0 if i == 0 else gap})
        if sb:
            lines.append({'t': ('    ' + sb), 'size': sub, 'color': MUTE, 'sb': 1, 'ls': 1.05})
    return text(s, l, t, w, Inches(4), lines)


# ============================================================================ SLIDE 0 — WELCOME / GREETING
s = new_slide(); grad(s, BG, GRAD2, 60)
chrome(s, 0); dotgrid(s, Inches(10.4), Inches(1.0), 8, 9)
try:
    s.shapes.add_picture(LOGO, SW - Inches(1.95), Inches(0.85), height=Inches(0.95))
except Exception: pass
anim_reset()
text(s, MARGIN, Inches(1.5), Inches(9), Inches(0.4),
     [{'t': 'THESIS DEFENSE  ·  WELCOME', 'size': 12, 'color': CYAN, 'bold': True}], track=2.5)
mark()
text(s, MARGIN, Inches(2.0), Inches(10.5), Inches(0.8),
     [{'t': 'Good morning everyone.', 'size': 36, 'color': WHITE, 'bold': True, 'font': TITLE}])
mark()
text(s, MARGIN, Inches(2.95), Inches(10.6), Inches(0.7),
     [{'t': 'First of all, I would like to greet the members of the committee and thank them for their presence.',
       'size': 15, 'color': TEXT, 'ls': 1.3}])
mark()
text(s, MARGIN, Inches(3.75), Inches(10.6), Inches(0.5),
     [{'t': 'Today, I am honored to present my graduation thesis, entitled:', 'size': 15, 'color': MUTE}])
mark()
card(s, MARGIN, Inches(4.4), Inches(11.0), Inches(1.15), fill=CARD2, line=CYAN_D, lw=1.2)
rect(s, MARGIN, Inches(4.4), Pt(3.5), Inches(1.15), fill=CYAN)
text(s, MARGIN + Inches(0.4), Inches(4.4), Inches(10.3), Inches(1.15),
     [{'t': 'Real-Time Shoplifting Detection in CCTV Videos\nUsing the YOLOv8 Framework',
       'size': 20, 'color': WHITE, 'bold': True, 'font': TITLE, 'ls': 1.15}], anchor=MIDDLE)
mark()
text(s, MARGIN, Inches(6.0), Inches(10), Inches(0.45),
     [{'t': 'Now, I will begin my presentation.', 'size': 14, 'color': CYAN, 'italic': True}])
mark()
finish(s, 'morph')
notes(s, "Defense opening — spoken greeting to the committee, thesis title in a highlighted card. Lines cascade in speaking order; Morph into the title slide.")

# ============================================================================ SLIDE 1
s = new_slide(); grad(s, BG, GRAD2, 60)
dotgrid(s, Inches(9.6), Inches(1.0), 9, 12)
corner_brackets(s, CYAN_D)
try:
    s.shapes.add_picture(LOGO, SW - Inches(1.95), Inches(0.95), height=Inches(0.95))
except Exception: pass
anim_reset()
text(s, MARGIN, Inches(1.7), Inches(9), Inches(0.4),
     [{'t': 'AI VISION SECURITY  ·  RETAIL LOSS PREVENTION', 'size': 12, 'color': CYAN, 'bold': True}], track=2.5)
mark()
text(s, MARGIN, Inches(2.2), Inches(10), Inches(1.7),
     [{'t': 'Theft Detection', 'size': 62, 'color': WHITE, 'bold': True, 'font': TITLE, 'sa': 6},
      {'t': 'AI that sees shoplifting before it happens.', 'size': 23, 'color': TEXT}])
mark()
rect(s, MARGIN, Inches(4.5), Inches(0.9), Pt(3), fill=CYAN)
mark()
text(s, MARGIN, Inches(5.0), Inches(10.5), Inches(0.9),
     [{'t': 'Real-time pose + object intelligence that watches every camera, flags suspicious behaviour, and alerts your team in seconds — across local webcams and RTSP networks.',
       'size': 14, 'color': MUTE, 'ls': 1.25}])
mark()
text(s, MARGIN, Inches(6.25), Inches(11), Inches(0.5),
     [{'runs': [{'t': 'Final-year project  ·  ', 'size': 12.5, 'color': MUTE},
                {'t': 'Université de Djelfa', 'size': 12.5, 'color': TEXT, 'bold': True},
                {'t': '   ·   Live demo: ', 'size': 12.5, 'color': MUTE},
                {'t': 'theft-detection-dusky.vercel.app', 'size': 12.5, 'color': CYAN, 'bold': True}]}])
finish(s, 'morph')
notes(s, "Hero / title. Transition: Morph. Build: kicker → wordmark+tagline → rule → descriptor → credits, auto fade-in cascade.")

# ============================================================================ SLIDE 3
s = new_slide(); header(s, "The Problem", "Theft is fast, silent — and human eyes miss it", 3)
text(s, MARGIN, Inches(2.4), Inches(11.4), Inches(0.7),
     [{'t': 'A single guard cannot watch dozens of feeds at once. Concealment happens in seconds, and traditional CCTV only proves what was lost — long after the thief has left.',
       'size': 14, 'color': MUTE, 'ls': 1.25}])
mark()
items = [
    ("Seconds, not minutes", "A pick-up-and-pocket gesture takes 2–3 seconds — far faster than a human scan of the floor."),
    ("Walls of blind video", "Dozens of CCTV streams, zero understanding. Nobody is watching the screen that matters."),
    ("No proof, no response", "Footage is reviewed after the loss; there is no live alert and rarely court-usable evidence."),
    ("Staff at risk", "Untrained confrontation is dangerous — teams need an early, objective signal, not a guess."),
]
y = Inches(3.35)
for i, (h, sub) in enumerate(items):
    rect(s, MARGIN, y + Inches(0.04), Pt(3), Inches(0.75), fill=CYAN)
    text(s, MARGIN + Inches(0.35), y, Inches(11.0), Inches(0.85),
         [{'t': h, 'size': 16.5, 'color': WHITE, 'bold': True, 'sa': 2},
          {'t': sub, 'size': 12.5, 'color': MUTE, 'ls': 1.1}])
    mark()
    y += Inches(0.95)
finish(s, 'fade')
notes(s, "Problem. Intro paragraph then 4 pain points with sub-detail, cascade fade-in.")

# ============================================================================ SLIDE 5
s = new_slide(); header(s, "Why Now", "Legacy CCTV records the past. It can't prevent it.", 5, tsize=28)
gap = Inches(0.5); cwd = (CW - gap) / 2
old = ["Passive recording only", "Reviewed after the loss", "One guard, many screens", "No alerts, no analytics", "Evidence often unusable"]
new = ["Real-time AI analysis", "Alerts the moment it happens", "Watches every feed at once", "Trends, heatmaps + CSV export", "Timestamped snapshot evidence"]
card(s, MARGIN, Inches(2.5), cwd, Inches(3.7), fill=CARD)
text(s, MARGIN + Inches(0.4), Inches(2.78), cwd - Inches(0.8), Inches(0.4), [{'t': 'TRADITIONAL CCTV', 'size': 12.5, 'color': MUTE, 'bold': True}], track=2)
yy = Inches(3.35)
for it in old:
    text(s, MARGIN + Inches(0.4), yy, cwd - Inches(0.8), Inches(0.45),
         [{'runs': [{'t': '✕   ', 'size': 14, 'color': RED, 'bold': True}, {'t': it, 'size': 13.5, 'color': TEXT}]}])
    yy += Inches(0.53)
mark()
x2 = MARGIN + cwd + gap
card(s, x2, Inches(2.5), cwd, Inches(3.7), fill=CARD2, line=CYAN, lw=1.5)
text(s, x2 + Inches(0.4), Inches(2.78), cwd - Inches(0.8), Inches(0.4), [{'t': 'THEFT DETECTION', 'size': 12.5, 'color': CYAN, 'bold': True}], track=2)
yy = Inches(3.35)
for it in new:
    text(s, x2 + Inches(0.4), yy, cwd - Inches(0.8), Inches(0.45),
         [{'runs': [{'t': '✓   ', 'size': 14, 'color': GREEN, 'bold': True}, {'t': it, 'size': 13.5, 'color': WHITE}]}])
    yy += Inches(0.53)
mark()
finish(s, 'fade')
notes(s, "Why now. Before/after compare, left then right card fade-in; right highlighted.")

# ============================================================================ SLIDE 7
s = new_slide(); header(s, "The Solution", "One platform, four ways it protects you", 7)
pill = [("01", "Detect", "Pose + object AI flags risky behaviour live.", ["Concealment & loitering", "ROI intrusion & posture"]),
        ("02", "Alert", "Instant, multi-channel response with evidence.", ["Browser siren on-screen", "Telegram + email snapshot"]),
        ("03", "Analyze", "Turn any recording into a forensic review.", ["Upload MP4, frame-by-frame", "Annotated video returned"]),
        ("04", "Recognize", "Match the faces that matter automatically.", ["Blacklist auto-alarm", "VIP / staff whitelist"])]
gap = Inches(0.3); cwd = (CW - 3 * gap) / 4
for i, (n, h, claim, bl) in enumerate(pill):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.5), cwd, Inches(3.7))
    chip(s, x + Inches(0.28), Inches(2.8), Inches(0.62), n, fill=BLUE)
    text(s, x + Inches(0.28), Inches(3.6), cwd - Inches(0.56), Inches(0.4), [{'t': h, 'size': 18, 'color': WHITE, 'bold': True}])
    text(s, x + Inches(0.28), Inches(4.05), cwd - Inches(0.56), Inches(0.9), [{'t': claim, 'size': 11.5, 'color': MUTE, 'ls': 1.15}])
    yy = Inches(5.0)
    for b in bl:
        text(s, x + Inches(0.28), yy, cwd - Inches(0.5), Inches(0.4),
             [{'runs': [{'t': '› ', 'size': 11.5, 'color': CYAN, 'bold': True}, {'t': b, 'size': 11, 'color': TEXT}]}])
        yy += Inches(0.42)
    mark()
finish(s, 'fade')
notes(s, "Solution pillars. 4 cards each claim + 2 bullets, cascade L→R. Verbs structure the deck.")

# ============================================================================ SLIDE 9
s = new_slide(); header(s, "How It Works", "Three steps, milliseconds apart", 9)
steps = [("Watch", "Every feed, in parallel", "A dedicated thread per camera captures frames with zero queueing lag.", "ThreadedCamera"),
         ("Understand", "Posture, gesture, zone", "YOLOv8 pose + object models read keypoints and track items frame to frame.", "YOLOv8 · OpenCV"),
         ("Alert", "The instant risk crosses", "Threshold breached → on-screen siren plus Telegram & email with a snapshot.", "Siren · Bot · SMTP")]
gap = Inches(0.5); cwd = (CW - 2 * gap) / 3
for i, (h, claim, sub, tech) in enumerate(steps):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.7), cwd, Inches(3.2))
    text(s, x + Inches(0.35), Inches(2.55), Inches(2), Inches(1), [{'t': f'0{i+1}', 'size': 44, 'color': CYAN_D, 'bold': True, 'font': TITLE}])
    text(s, x + Inches(0.35), Inches(3.65), cwd - Inches(0.7), Inches(0.45), [{'t': h, 'size': 20, 'color': WHITE, 'bold': True}])
    text(s, x + Inches(0.35), Inches(4.1), cwd - Inches(0.7), Inches(0.4), [{'t': claim, 'size': 12.5, 'color': CYAN}])
    text(s, x + Inches(0.35), Inches(4.5), cwd - Inches(0.7), Inches(1.0), [{'t': sub, 'size': 11.5, 'color': MUTE, 'ls': 1.2}])
    text(s, x + Inches(0.35), Inches(5.5), cwd - Inches(0.7), Inches(0.35), [{'t': tech, 'size': 10.5, 'color': MUTE, 'bold': True, 'track': 1}])
    if i < 2:
        arrow(s, x + cwd + Inches(0.06), Inches(3.95), Inches(0.4), Inches(0.5))
    mark()
finish(s, 'fade')
notes(s, "Pipeline stepper. Each step card + chevron cascade; tech tag per step.")

# ============================================================================ SLIDE 10
s = new_slide(); header(s, "Architecture", "A clean pipeline from lens to alert", 10)
flow = [("CAMERAS", "Webcam / RTSP", "Threaded multi-feed capture"),
        ("AI ENGINE", "YOLOv8 Pose + Object", "Per-(cam, track) state logic"),
        ("DASHBOARD", "Next.js + WebSocket", "Live view, ROI, analytics"),
        ("ALERTS", "Siren · Telegram · Email", "Snapshot + SQLite log")]
gap = Inches(0.55); bw = (CW - 3 * gap) / 4
for i, (h, sub, det) in enumerate(flow):
    x = MARGIN + i * (bw + gap)
    fill = CARD2 if i == 1 else CARD
    ln = CYAN if i == 1 else LINE
    card(s, x, Inches(2.9), bw, Inches(2.2), fill=fill, line=ln, lw=1.5 if i == 1 else 1.0)
    text(s, x, Inches(3.15), bw, Inches(0.4), [{'t': h, 'size': 12.5, 'color': CYAN, 'bold': True, 'align': CENTER}], align=CENTER, track=1.5)
    text(s, x + Inches(0.2), Inches(3.65), bw - Inches(0.4), Inches(0.5), [{'t': sub, 'size': 12.5, 'color': WHITE, 'bold': True, 'align': CENTER, 'ls': 1.05}], align=CENTER)
    text(s, x + Inches(0.2), Inches(4.35), bw - Inches(0.4), Inches(0.6), [{'t': det, 'size': 10.5, 'color': MUTE, 'align': CENTER, 'ls': 1.1}], align=CENTER)
    if i < 3:
        arrow(s, x + bw + Inches(0.1), Inches(3.7), Inches(0.34), Inches(0.55))
    mark()
text(s, MARGIN, Inches(5.5), CW, Inches(0.8),
     [{'t': 'FastAPI serves async endpoints and a live WebSocket stream; SQLite persists every alert and face encoding; psutil reports CPU/RAM to the dashboard in real time.',
       'size': 12.5, 'color': MUTE, 'align': CENTER, 'ls': 1.25}], align=CENTER)
mark()
finish(s, 'fade')
notes(s, "Architecture flow with per-stage tech sub-lines + infra line. Boxes cascade.")

# ============================================================================ SLIDE 11
s = new_slide(); header(s, "Capability · Performance", "Multi-threaded camera engine", 11)
bullets(s, MARGIN, Inches(2.55), Inches(6.2), [
    ("Zero-lag parallel capture", "Each camera runs on its own thread — frames never queue behind a slow feed."),
    ("No identity bleed", "Tracker state is isolated per (camera, person) pair, even with crowds in frame."),
    ("Resilient by design", "A stalled or dropped stream can't freeze the rest of the system."),
    ("Resource-aware", "Object detection runs every Nth frame to balance accuracy against load."),
])
mark()
metric(s, MARGIN, Inches(5.55), Inches(2.9), "Real-time", "live inference per feed")
metric(s, MARGIN + Inches(3.15), Inches(5.55), Inches(2.9), "N× feeds", "concurrent, no lock-step")
mark()
bx = Inches(7.7)
for i in range(4):
    yy = Inches(2.6) + i * Inches(0.72)
    card(s, bx, yy, Inches(2.9), Inches(0.56), fill=CARD, line=LINE)
    text(s, bx + Inches(0.22), yy, Inches(2.5), Inches(0.56), [{'t': f'CAM 0{i+1}', 'size': 11.5, 'color': MUTE, 'bold': True}], anchor=MIDDLE)
    rect(s, bx + Inches(2.0), yy + Inches(0.2), Inches(0.65), Inches(0.16), fill=CYAN if i % 2 == 0 else BLUE)
    arrow(s, bx + Inches(2.95), yy + Inches(0.1), Inches(0.32), Inches(0.34), color=LINE)
    mark()
card(s, bx + Inches(3.4), Inches(3.35), Inches(1.55), Inches(1.25), fill=CARD2, line=CYAN, lw=1.5)
text(s, bx + Inches(3.4), Inches(3.55), Inches(1.55), Inches(0.9), [{'t': 'AI\nCORE', 'size': 15, 'color': WHITE, 'bold': True, 'align': CENTER}], align=CENTER, anchor=MIDDLE)
mark()
finish(s, 'fade')
notes(s, "Engine. Denser bullets + 2 metric chips + schematic that cascades into the AI core.")

# ============================================================================ SLIDE 12
s = new_slide(); header(s, "Capability · Intelligence", "It reads behaviour, not just pixels", 12)
feats = [("Concealment", "Item hidden on the body", ["Item pick-up detected", "Hand → pocket / bag gesture"]),
         ("Loitering", "Dwell beyond a threshold", ["Per-zone time limit", "Configurable seconds"]),
         ("Zone intrusion", "Restricted-area breach", ["Wrist crosses the ROI line", "Instant siren trigger"]),
         ("Suspicious posture", "Unusual body signals", ["Sudden bending detected", "Low-visibility aisle watch"])]
gap = Inches(0.4); cwd = (CW - gap) / 2; ch = Inches(1.78)
for i, (h, claim, bl) in enumerate(feats):
    col, row = i % 2, i // 2
    x = MARGIN + col * (cwd + gap); y = Inches(2.5) + row * (ch + Inches(0.3))
    card(s, x, y, cwd, ch)
    chip(s, x + Inches(0.3), y + Inches(0.3), Inches(0.5), str(i + 1), fill=BLUE, fsize=13)
    text(s, x + Inches(1.0), y + Inches(0.26), cwd - Inches(1.3), Inches(0.4), [{'t': h, 'size': 17, 'color': WHITE, 'bold': True}])
    text(s, x + Inches(1.0), y + Inches(0.66), cwd - Inches(1.3), Inches(0.4), [{'t': claim, 'size': 11.5, 'color': CYAN}])
    yy = y + Inches(1.02)
    for b in bl:
        text(s, x + Inches(1.0), yy, cwd - Inches(1.2), Inches(0.35),
             [{'runs': [{'t': '› ', 'size': 11, 'color': CYAN, 'bold': True}, {'t': b, 'size': 11, 'color': MUTE}]}])
        yy += Inches(0.34)
    mark()
finish(s, 'fade')
notes(s, "2×2 behaviour grid, each cell claim + 2 bullets. Z-pattern cascade.")

# ============================================================================ SLIDE 15
s = new_slide(); header(s, "Capability · Identity", "Faces that matter, matched instantly", 15)
gap = Inches(0.5); cwd = (CW - gap) / 2
card(s, MARGIN, Inches(2.5), cwd, Inches(3.6), fill=CARD, line=RED, lw=1.3)
chip(s, MARGIN + Inches(0.4), Inches(2.8), Inches(0.65), '!', fill=RED, fsize=22)
text(s, MARGIN + Inches(1.2), Inches(2.9), cwd - Inches(1.5), Inches(0.5), [{'t': 'Blacklist', 'size': 21, 'color': WHITE, 'bold': True}])
text(s, MARGIN + Inches(1.2), Inches(3.35), cwd - Inches(1.5), Inches(0.4), [{'t': 'Known offenders, flagged on sight', 'size': 12, 'color': RED}])
for k, b in enumerate(["High-priority alarm on entry", "Evidence auto-recorded", "Dlib CNN face encodings in SQLite"]):
    text(s, MARGIN + Inches(0.4), Inches(4.05) + k * Inches(0.52), cwd - Inches(0.8), Inches(0.45),
         [{'runs': [{'t': '› ', 'size': 12.5, 'color': RED, 'bold': True}, {'t': b, 'size': 12.5, 'color': TEXT}]}])
mark()
x2 = MARGIN + cwd + gap
card(s, x2, Inches(2.5), cwd, Inches(3.6), fill=CARD, line=GREEN, lw=1.3)
chip(s, x2 + Inches(0.4), Inches(2.8), Inches(0.65), '★', fill=GREEN, fsize=19)
text(s, x2 + Inches(1.2), Inches(2.9), cwd - Inches(1.5), Inches(0.5), [{'t': 'VIP Whitelist', 'size': 21, 'color': WHITE, 'bold': True}])
text(s, x2 + Inches(1.2), Inches(3.35), cwd - Inches(1.5), Inches(0.4), [{'t': 'Trusted staff & loyal customers', 'size': 12, 'color': GREEN}])
for k, b in enumerate(["Recognised, not alarmed", "Green greeting badge shown", "One-click register & delete"]):
    text(s, x2 + Inches(0.4), Inches(4.05) + k * Inches(0.52), cwd - Inches(0.8), Inches(0.45),
         [{'runs': [{'t': '› ', 'size': 12.5, 'color': GREEN, 'bold': True}, {'t': b, 'size': 12.5, 'color': TEXT}]}])
mark()
finish(s, 'fade')
notes(s, "Identity. Blacklist vs VIP cards, each with bullets + how it works.")

# ============================================================================ SLIDE 16
s = new_slide(); header(s, "Capability · Response", "The alarm reaches you anywhere", 16)
chans = [("Browser siren", "On-screen, instant", ["Web Audio synth — no MP3 lag", "Sweeping emergency tone"]),
         ("Telegram", "On your phone", ["Bot pushes alert + snapshot", "Delivered in seconds"]),
         ("Email (SMTP)", "For the record", ["Automated mail with evidence", "Shared with the whole team"])]
gap = Inches(0.4); cwd = (CW - 2 * gap) / 3
for i, (h, claim, bl) in enumerate(chans):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.6), cwd, Inches(3.1))
    chip(s, x + Inches(0.35), Inches(2.9), Inches(0.6), ['♪', '✈', '@'][i], fill=BLUE, fsize=17)
    text(s, x + Inches(0.35), Inches(3.7), cwd - Inches(0.7), Inches(0.4), [{'t': h, 'size': 17, 'color': WHITE, 'bold': True}])
    text(s, x + Inches(0.35), Inches(4.1), cwd - Inches(0.7), Inches(0.35), [{'t': claim, 'size': 11.5, 'color': CYAN}])
    yy = Inches(4.55)
    for b in bl:
        text(s, x + Inches(0.35), yy, cwd - Inches(0.6), Inches(0.4),
             [{'runs': [{'t': '› ', 'size': 11.5, 'color': CYAN, 'bold': True}, {'t': b, 'size': 11, 'color': MUTE}]}])
        yy += Inches(0.4)
    mark()
text(s, MARGIN, Inches(6.0), CW, Inches(0.4),
     [{'t': 'A smart 3-second cooldown prevents alert fatigue.', 'size': 12.5, 'color': CYAN, 'align': CENTER, 'bold': True}], align=CENTER)
mark()
finish(s, 'fade')
notes(s, "Response triptych, each channel claim + 2 bullets; cooldown line last.")

# ============================================================================ SLIDE 17
s = new_slide(); header(s, "Capability · Forensics", "Upload a video. Get answers.", 17)
bullets(s, MARGIN, Inches(2.5), Inches(3.95), [
    ("Any MP4, fully analysed", "The same live AI runs frame-by-frame over recorded footage."),
    ("Annotated video returned", "Download the clip with detections drawn on every frame."),
    ("Every alert listed", "A clear summary of what was flagged, and when."),
    ("Telegram on detection", "Optional push notification the moment something is found."),
])
mark()
shot(s, f"{SHOTS}/video-analysis.png", Inches(5.15), Inches(2.0), Inches(7.3))
mark()
finish(s, 'fade')
notes(s, "Forensics. Denser bullets left; product screenshot Float-in right. Demo live if possible.")

# ============================================================================ SLIDES 18-22 screenshots
def shot_slide(n, kicker, title, img, caption_items):
    s = new_slide(); header(s, kicker, title, n, tsize=29)
    bullets(s, MARGIN, Inches(2.5), Inches(3.75), caption_items); mark()
    shot(s, f"{SHOTS}/{img}", Inches(5.0), Inches(2.1), Inches(7.45)); mark()
    finish(s, 'fade')
    notes(s, f"Product · {title}. Screenshot Float-in + feature bullets cascade.")

shot_slide(18, "The Product", "One glassmorphic control room", "overview.png", [
    ("Live multi-camera wall", "Every feed streams into one view over a real-time WebSocket."),
    ("System telemetry", "Live CPU and RAM bars, fed by psutil from the backend."),
    ("Weekly trends", "Suspicious vs. reviewed events, charted by Recharts."),
    ("Today's alert curve", "Hour-by-hour activity, visible at a glance."),
])
shot_slide(19, "The Product", "Camera setup & ROI drawer", "cameras.png", [
    ("Add any source", "Local webcams or RTSP network cameras in seconds."),
    ("Draw security zones", "Polygon ROIs over the live feed, per camera."),
    ("Saved persistently", "Definitions and coordinates stored in cameras.json."),
    ("Resolution-safe", "Zones map exactly to the 1280×720 matrix."),
])
shot_slide(20, "The Product", "Face management panel", "faces.png", [
    ("Upload portraits", "Register new faces with a single photo."),
    ("Tag the role", "Blacklist for offenders, VIP for trusted people."),
    ("Instant DB sync", "One-click SQLite delete with live memory update."),
    ("Graceful fallback", "Runs fine even where face libs aren't installed."),
])
shot_slide(21, "The Product", "Alerts, trends & CSV export", "history.png", [
    ("Search & filter", "By ID, message, camera or event category."),
    ("Weekly trends", "Suspicious vs reviewed events, visualised."),
    ("One-click export", "Download filtered alerts as a clean CSV."),
    ("Evidence on file", "Every alert keeps its timestamped snapshot."),
])
shot_slide(22, "The Product", "Notifications setup", "settings.png", [
    ("Telegram in minutes", "Paste bot token + chat ID, send a test."),
    ("SMTP email alerts", "Automated mail with snapshot attached."),
    ("Broadcast instantly", "Push alerts to every configured channel."),
    ("Tune the behaviour", "Cooldowns and thresholds in one place."),
])

# ============================================================================ SLIDE 23
s = new_slide(); header(s, "Under the Hood", "Built on a proven, modern stack", 23)
groups = [("AI / VISION", ["YOLOv8 Pose", "YOLOv8 Object", "OpenCV", "Dlib face encodings"], "Real-time detection & tracking"),
          ("BACKEND", ["Python 3.10+", "FastAPI", "WebSockets", "SQLite"], "Async API + live streaming"),
          ("FRONTEND", ["Next.js 16", "React 19", "Tailwind CSS", "Recharts"], "Glassmorphic control room")]
gap = Inches(0.4); cwd = (CW - 2 * gap) / 3
for i, (h, items, why) in enumerate(groups):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.5), cwd, Inches(3.5))
    text(s, x + Inches(0.4), Inches(2.78), cwd - Inches(0.8), Inches(0.4), [{'t': h, 'size': 12.5, 'color': CYAN, 'bold': True}], track=2)
    yy = Inches(3.35)
    for it in items:
        oval(s, x + Inches(0.4), yy + Inches(0.08), Inches(0.11), Inches(0.11), CYAN)
        text(s, x + Inches(0.68), yy, cwd - Inches(1.0), Inches(0.42), [{'t': it, 'size': 14, 'color': TEXT}])
        yy += Inches(0.5)
    text(s, x + Inches(0.4), Inches(5.45), cwd - Inches(0.8), Inches(0.45), [{'t': why, 'size': 11, 'color': MUTE, 'italic': True}])
    mark()
finish(s, 'fade')
notes(s, "Tech stack, 3 grouped columns + one-line rationale each. Columns cascade.")

# ============================================================================ SLIDE 25
s = new_slide(); header(s, "Differentiation", "Why we win", 25)
rows = [["", "Legacy CCTV", "Cloud SaaS", "Theft Detection"],
        ["Real-time prevention", "✕", "~", "✓"],
        ["Behaviour + pose AI", "✕", "~", "✓"],
        ["Runs fully on-premise", "✓", "✕", "✓"],
        ["Video forensic upload", "✕", "~", "✓"],
        ["One-time, no per-seat fee", "✓", "✕", "✓"]]
tx, ty = MARGIN, Inches(2.35); tw = CW; rh = Inches(0.58); c0 = Inches(4.6); cw = (tw - c0) / 3
for r, row in enumerate(rows):
    y = ty + r * rh; head = (r == 0)
    if head: rect(s, tx, y, tw, rh, fill=CARD2)
    elif r % 2 == 1: rect(s, tx, y, tw, rh, fill=CARD)
    text(s, tx + Inches(0.3), y, c0 - Inches(0.4), rh, [{'t': row[0], 'size': 13.5, 'color': WHITE if head else TEXT, 'bold': head}], anchor=MIDDLE)
    for c in range(1, 4):
        cx = tx + c0 + (c - 1) * cw; val = row[c]
        if c == 3:
            rect(s, cx, y, cw, rh, fill=None, line=CYAN, lw=1.0) if head else rect(s, cx, y, cw, rh, fill=CARD2)
        col = WHITE
        if val == '✓': col = GREEN
        elif val == '✕': col = RED
        elif val == '~': col = AMBER
        text(s, cx, y, cw, rh, [{'t': val, 'size': 15 if not head else 12.5, 'color': CYAN if head and c == 3 else (col if not head else MUTE), 'bold': True, 'align': CENTER}], align=CENTER, anchor=MIDDLE)
    mark()
rect(s, tx, ty, tw, rh * len(rows), fill=None, line=LINE, lw=1.0)
text(s, MARGIN, Inches(6.15), CW, Inches(0.4), [{'t': 'On-premise control + real-time prevention + a forensic upload tool — without per-seat cloud fees.   ✓ full · ~ partial · ✕ none',
     'size': 11.5, 'color': MUTE, 'align': CENTER}], align=CENTER); mark()
finish(s, 'fade')
notes(s, "Comparison matrix, our column highlighted; rows wipe in; takeaway + legend last.")

# ============================================================================ SLIDE 27
s = new_slide(); header(s, "Business Model", "Simple, scalable pricing", 27)
tiers = [("Starter", "Single store", ["1 location", "Up to 4 cameras", "Core detection", "Email alerts"], False),
         ("Pro", "Growing retail", ["Up to 5 locations", "Unlimited cameras", "Face recognition", "Telegram + analytics"], True),
         ("Enterprise", "Chains & malls", ["Unlimited locations", "On-prem / private cloud", "Custom AI model", "Priority support"], False)]
gap = Inches(0.4); cwd = (CW - 2 * gap) / 3
for i, (name, who, feats, hi) in enumerate(tiers):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.45), cwd, Inches(3.75), fill=CARD2 if hi else CARD, line=CYAN if hi else LINE, lw=1.8 if hi else 1.0)
    if hi:
        rect(s, x + cwd / 2 - Inches(0.8), Inches(2.28), Inches(1.6), Inches(0.32), fill=CYAN, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        text(s, x + cwd / 2 - Inches(0.8), Inches(2.28), Inches(1.6), Inches(0.32), [{'t': 'POPULAR', 'size': 9.5, 'color': BG, 'bold': True, 'align': CENTER}], align=CENTER, anchor=MIDDLE, track=1.5)
    text(s, x + Inches(0.35), Inches(2.72), cwd - Inches(0.7), Inches(0.5), [{'t': name, 'size': 21, 'color': WHITE, 'bold': True}])
    text(s, x + Inches(0.35), Inches(3.2), cwd - Inches(0.7), Inches(0.4), [{'t': who, 'size': 12.5, 'color': CYAN}])
    yy = Inches(3.75)
    for f in feats:
        text(s, x + Inches(0.35), yy, cwd - Inches(0.7), Inches(0.45),
             [{'runs': [{'t': '✓  ', 'size': 12.5, 'color': GREEN, 'bold': True}, {'t': f, 'size': 13, 'color': TEXT}]}])
        yy += Inches(0.5)
    mark()
finish(s, 'fade')
notes(s, "Pricing tiers, middle highlighted with POPULAR. Maps to the BMC in the report.")


# --------------------------------------------------------------------- financial-table helper
def finrow(s, x, y, w, widths, cells, style='norm', rh=Inches(0.40), zebra=False):
    if style in ('head', 'sub'):
        rect(s, x, y, w, rh, fill=CARD2)
    elif style == 'tot':
        rect(s, x, y, w, rh, fill=CARD2, line=CYAN, lw=1.2)
    elif zebra:
        rect(s, x, y, w, rh, fill=CARD)
    cx = x
    n = len(cells)
    for i, cell in enumerate(cells):
        cw_ = int(w * widths[i])
        al = RIGHT if i == n - 1 else LEFT
        if style == 'head':
            spec = {'t': cell, 'size': 10, 'color': MUTE, 'bold': True, 'align': al, 'track': 1}
        elif style == 'sub':
            spec = {'t': cell, 'size': 11.5, 'color': WHITE, 'bold': True, 'align': al}
        elif style == 'tot':
            spec = {'t': cell, 'size': 12, 'color': CYAN if i == n - 1 else WHITE, 'bold': True, 'align': al}
        else:
            spec = {'t': cell, 'size': 11.5, 'color': TEXT if i else MUTE, 'align': al}
        pad = Inches(0.16)
        text(s, cx + pad, y, cw_ - 2 * pad, rh, [spec], align=al, anchor=MIDDLE)
        cx += cw_


def barchart(s, x0, base, chart_w, maxh, vals, vmax, months, color_fn, label_y_off=Inches(0.10)):
    slot = int(chart_w) // len(vals)
    bw = int(slot * 0.66)
    for i, v in enumerate(vals):
        bh = int(int(maxh) * abs(v) / vmax)
        bx = int(x0) + i * slot + (slot - bw) // 2
        if v >= 0:
            rect(s, bx, int(base) - bh, bw, max(bh, 1), fill=color_fn(i, v))
        else:
            rect(s, bx, int(base), bw, max(bh, 1), fill=color_fn(i, v))
        text(s, int(x0) + i * slot, int(base) + int(label_y_off), slot, Inches(0.3),
             [{'t': months[i], 'size': 7.5, 'color': MUTE, 'align': CENTER}], align=CENTER)
        mark()
    rect(s, int(x0), int(base), int(chart_w), Pt(1), fill=LINE)


# ============================================================================ SLIDE F1 — ASSETS (3.1)
s = new_slide(); header(s, "Financial Study · Report Ch. 3", "Assets and liabilities", 0)
tw = Inches(7.0); widths = [0.22, 0.50, 0.28]
finrow(s, MARGIN, Inches(2.35), tw, widths, ["TYPE", "ITEM", "VALUE (DZD)"], style='head'); mark()
rows = [("Financial", "Cash", "0"), ("Financial", "Inventory — licences & spare cameras", "250,000"),
        ("Fixed", "Workstation / GPU server", "350,000"), ("Fixed", "IP / CCTV cameras (demo kit)", "180,000"),
        ("Fixed", "Network & storage equipment", "90,000"), ("Fixed", "Office equipment", "60,000"),
        ("Fixed", "Outdoor sign", "10,000"), ("Fixed", "Decorations & furniture", "40,000")]
yy = Inches(2.75)
for k, r in enumerate(rows):
    finrow(s, MARGIN, yy, tw, widths, list(r), zebra=(k % 2 == 0)); mark()
    yy += Inches(0.40)
finrow(s, MARGIN, yy, tw, widths, ["", "Total assets", "980,000"], style='tot'); mark()
px = MARGIN + tw + Inches(0.45); pw = SW - MARGIN - px
card(s, px, Inches(2.35), pw, Inches(2.0), fill=CARD2, line=CYAN_D, lw=1.1)
text(s, px + Inches(0.35), Inches(2.65), pw - Inches(0.7), Inches(0.9),
     [{'t': '980K', 'size': 40, 'color': CYAN, 'bold': True, 'font': TITLE}])
text(s, px + Inches(0.35), Inches(3.6), pw - Inches(0.7), Inches(0.6),
     [{'t': 'DZD opening assets — an asset-light venture', 'size': 12, 'color': WHITE, 'bold': True, 'ls': 1.15}])
mark()
for k, b in enumerate(["GPU workstation is the main investment", "No costly cloud infrastructure needed", "Demo camera kit doubles as install stock"]):
    text(s, px, Inches(4.65) + k * Inches(0.55), pw, Inches(0.5),
         [{'runs': [{'t': '› ', 'size': 12, 'color': CYAN, 'bold': True}, {'t': b, 'size': 12, 'color': TEXT}]}])
    mark()
finish(s, 'fade')
notes(s, "Financial study 3.1 — opening assets table (Total 980,000 DZD) + asset-light takeaway panel. Rows cascade.")

# ============================================================================ SLIDE F2 — REVENUES (3.2)
s = new_slide(); header(s, "Financial Study · Report Ch. 3", "Revenues — first-year projection", 0)
tw = Inches(5.6); widths = [0.62, 0.38]
finrow(s, MARGIN, Inches(2.45), tw, widths, ["REVENUE STREAM", "YEAR (DZD)"], style='head'); mark()
rows = [("Subscriptions (per camera / site)", "4,180,000"), ("Installation & setup", "1,620,000"), ("Support & training", "200,000")]
yy = Inches(2.85)
for k, r in enumerate(rows):
    finrow(s, MARGIN, yy, tw, widths, list(r), zebra=(k % 2 == 0)); mark()
    yy += Inches(0.40)
finrow(s, MARGIN, yy, tw, widths, ["Total first-year revenue", "6,000,000"], style='tot'); mark()
text(s, MARGIN, yy + Inches(0.65), tw, Inches(1.5),
     [{'t': 'Recurring subscriptions are ~70% of income and compound as new shops are onboarded — installation fees are one-off, support grows with the installed base.',
       'size': 12.5, 'color': MUTE, 'ls': 1.3}])
mark()
text(s, Inches(7.05), Inches(2.45), Inches(5.35), Inches(0.4),
     [{'t': 'MONTHLY TOTAL REVENUE (K DZD)', 'size': 10.5, 'color': CYAN, 'bold': True, 'track': 1.5}])
mark()
barchart(s, Inches(7.05), Inches(5.95), Inches(5.35), Inches(2.6),
         [180, 200, 300, 340, 395, 455, 530, 550, 625, 665, 760, 800], 800,
         list('JFMAMJJASOND'), lambda i, v: CYAN if i >= 10 else BLUE)
text(s, Inches(7.05), Inches(6.35), Inches(5.35), Inches(0.35),
     [{'t': '180K in January  →  800K in December  (×4.4 growth)', 'size': 10.5, 'color': MUTE, 'align': CENTER}], align=CENTER)
mark()
finish(s, 'fade')
notes(s, "Financial study 3.2 — revenue streams table (Total 6.0M DZD) + monthly growth bar chart, bars cascade L→R.")

# ============================================================================ SLIDE F3 — EXPENSES (3.3)
s = new_slide(); header(s, "Financial Study · Report Ch. 3", "Expenses — fixed and variable", 0)
tw = Inches(6.9); widths = [0.18, 0.54, 0.28]
finrow(s, MARGIN, Inches(2.30), tw, widths, ["TYPE", "DESCRIPTION", "YEAR (DZD)"], style='head', rh=Inches(0.38)); mark()
rows = [("Fixed", "Rent", "360,000", 'norm'), ("Fixed", "Salaries", "960,000", 'norm'),
        ("", "Total fixed  ·  110K / month, flat", "1,320,000", 'sub'),
        ("Variable", "Engineers (contract)", "1,010,000", 'norm'), ("Variable", "Hardware / cameras", "540,000", 'norm'),
        ("Variable", "Marketing & advertising", "295,000", 'norm'), ("Variable", "Cloud / hosting", "72,000", 'norm'),
        ("Variable", "Other expenses", "60,000", 'norm'),
        ("", "Total variable  ·  scales with installs", "1,977,000", 'sub')]
yy = Inches(2.68)
for k, (a, b, c, st) in enumerate(rows):
    finrow(s, MARGIN, yy, tw, widths, [a, b, c], style=st, rh=Inches(0.38), zebra=(st == 'norm' and k % 2 == 0)); mark()
    yy += Inches(0.38)
finrow(s, MARGIN, yy, tw, widths, ["", "Total costs", "3,297,000"], style='tot', rh=Inches(0.40)); mark()
px = MARGIN + tw + Inches(0.45); pw = SW - MARGIN - px
card(s, px, Inches(2.30), pw, Inches(1.7), fill=CARD2, line=CYAN_D, lw=1.1)
text(s, px + Inches(0.3), Inches(2.55), pw - Inches(0.6), Inches(0.8),
     [{'t': '3.3M', 'size': 34, 'color': AMBER, 'bold': True, 'font': TITLE}])
text(s, px + Inches(0.3), Inches(3.35), pw - Inches(0.6), Inches(0.5),
     [{'t': 'DZD total first-year costs', 'size': 11.5, 'color': WHITE, 'bold': True}])
mark()
# fixed vs variable split bar
text(s, px, Inches(4.35), pw, Inches(0.3), [{'t': 'COST SPLIT', 'size': 10, 'color': CYAN, 'bold': True, 'track': 1.5}])
rect(s, px, Inches(4.7), pw, Inches(0.34), fill=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
rect(s, px, Inches(4.7), int(pw * 0.40), Inches(0.34), fill=BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
text(s, px, Inches(5.12), pw, Inches(0.35),
     [{'runs': [{'t': '40% fixed', 'size': 10.5, 'color': BLUE, 'bold': True}, {'t': '   ·   ', 'size': 10.5, 'color': MUTE}, {'t': '60% variable', 'size': 10.5, 'color': MUTE, 'bold': True}]}])
mark()
text(s, px, Inches(5.7), pw, Inches(1.0),
     [{'t': 'Variable spend follows the install pipeline — costs only grow when revenue does.', 'size': 11.5, 'color': MUTE, 'ls': 1.25}])
mark()
finish(s, 'fade')
notes(s, "Financial study 3.3 — expenses table (Total 3,297,000 DZD) + cost-split panel. Rows cascade.")

# ============================================================================ SLIDE F4 — CASH FLOWS (3.4)
s = new_slide(); header(s, "Financial Study · Report Ch. 3", "Cash flows — from deficit to surplus", 0)
tw = Inches(4.7); widths = [0.60, 0.40]
finrow(s, MARGIN, Inches(2.45), tw, widths, ["FLOW", "YEAR (DZD)"], style='head'); mark()
finrow(s, MARGIN, Inches(2.85), tw, widths, ["Total inflows", "6,000,000"], zebra=True); mark()
finrow(s, MARGIN, Inches(3.25), tw, widths, ["Total outflows", "3,297,000"]); mark()
finrow(s, MARGIN, Inches(3.65), tw, widths, ["Net cash flow", "+2,703,000"], style='tot'); mark()
for k, b in enumerate([("January only is negative", "−16K during initial setup"), ("Break-even from month 2", "+4K in February, rising every month"), ("Net margin ≈ 45%", "of first-year revenue retained")]):
    text(s, MARGIN, Inches(4.55) + k * Inches(0.62), tw, Inches(0.6),
         [{'runs': [{'t': '› ', 'size': 12.5, 'color': CYAN, 'bold': True}, {'t': b[0] + ' — ', 'size': 12.5, 'color': WHITE, 'bold': True}, {'t': b[1], 'size': 12, 'color': MUTE}]}])
    mark()
text(s, Inches(6.2), Inches(2.45), Inches(6.2), Inches(0.4),
     [{'t': 'MONTHLY NET CASH FLOW (K DZD)', 'size': 10.5, 'color': CYAN, 'bold': True, 'track': 1.5}])
mark()
barchart(s, Inches(6.2), Inches(5.45), Inches(6.2), Inches(2.45),
         [-16, 4, 69, 109, 154, 179, 244, 264, 304, 344, 404, 444], 444,
         list('JFMAMJJASOND'), lambda i, v: RED if v < 0 else (GREEN if v > 300 else CYAN),
         label_y_off=Inches(0.18))
text(s, Inches(6.2), Inches(6.15), Inches(6.2), Inches(0.35),
     [{'t': 'a small setup deficit, then steady compounding surplus all year', 'size': 10.5, 'color': MUTE, 'align': CENTER}], align=CENTER)
mark()
finish(s, 'fade')
notes(s, "Financial study 3.4 — cash-flow summary (Net +2,703,000 DZD) + monthly net bar chart (Jan negative in red).")

# ============================================================================ SLIDE 28
s = new_slide(); header(s, "Financial Study · Summary", "The numbers behind the build", 28)
kpis = [("980K", "Assets (DZD)", CYAN), ("6.0M", "Revenue (DZD)", GREEN), ("3.3M", "Total cost (DZD)", AMBER), ("2.7M", "Net cash flow (DZD)", WHITE)]
gap = Inches(0.35); cwd = (CW - 3 * gap) / 4
for i, (num, lab, col) in enumerate(kpis):
    x = MARGIN + i * (cwd + gap)
    card(s, x, Inches(2.6), cwd, Inches(2.3))
    text(s, x, Inches(2.95), cwd, Inches(1.0), [{'t': num, 'size': 38, 'color': col, 'bold': True, 'font': TITLE, 'align': CENTER}], align=CENTER)
    text(s, x + Inches(0.2), Inches(4.05), cwd - Inches(0.4), Inches(0.6), [{'t': lab, 'size': 12, 'color': MUTE, 'align': CENTER}], align=CENTER)
    mark()
notes_items = [("Assumptions", "first operating cycle, single deployment region"), ("Revenue mix", "tiered subscriptions + enterprise licences"), ("Net positive", "cash flow turns positive within the cycle")]
yy = Inches(5.25)
for ld, sb in notes_items:
    text(s, MARGIN, yy, CW, Inches(0.4),
         [{'runs': [{'t': '› ', 'size': 12, 'color': CYAN, 'bold': True}, {'t': ld + ' — ', 'size': 12, 'color': WHITE, 'bold': True}, {'t': sb, 'size': 11.5, 'color': MUTE}]}])
    yy += Inches(0.42)
mark()
text(s, MARGIN, Inches(6.55), CW, Inches(0.3), [{'t': 'Consolidated from Tables 3.1 – 3.4 of the financial study.', 'size': 9.5, 'color': MUTE}]); mark()
finish(s, 'fade')
notes(s, "Financials. 4 KPI cards cascade + assumptions bullets. Net-positive is the headline.")

# ============================================================================ SLIDE BMC — BUSINESS MODEL CANVAS (Ch. 4)
s = new_slide(); header(s, "Report Ch. 4", "Business Model Canvas", 0)
text(s, MARGIN, Inches(2.18), CW, Inches(0.34),
     [{'t': 'THEFT DETECTION', 'size': 15, 'color': CYAN, 'bold': True, 'font': TITLE, 'align': CENTER, 'track': 6}], align=CENTER)
mark()


def bmc_cell(x, y, w, h, title, items, tsz=9.5, isz=8.5, inline=False):
    card(s, x, y, w, h, fill=CARD, line=LINE, radius=0.06)
    text(s, x + Inches(0.14), y + Inches(0.09), w - Inches(0.28), Inches(0.3),
         [{'t': title.upper(), 'size': tsz, 'color': CYAN, 'bold': True, 'track': 0.5}])
    if inline:
        text(s, x + Inches(0.14), y + Inches(0.42), w - Inches(0.28), h - Inches(0.5),
             [{'t': '   ·   '.join(items), 'size': isz, 'color': TEXT, 'ls': 1.1}])
    else:
        lines = [{'runs': [{'t': '› ', 'size': isz, 'color': CYAN_D, 'bold': True},
                           {'t': it, 'size': isz, 'color': TEXT}], 'sb': 0 if k == 0 else 2, 'ls': 1.0}
                 for k, it in enumerate(items)]
        text(s, x + Inches(0.14), y + Inches(0.40), w - Inches(0.28), h - Inches(0.5), lines)
    mark()


gx, gy = MARGIN, Inches(2.58)
ggap = Inches(0.10)
gcw = (CW - 4 * ggap) / 5
rh2 = Inches(1.46)                      # half-height cell
fh2 = rh2 * 2 + ggap                    # full-height cell
bot_y = gy + fh2 + ggap                 # bottom row y
bot_h = Inches(1.12)
c = lambda i: gx + i * (gcw + ggap)
bmc_cell(c(0), gy, gcw, fh2, "Key Partners",
         ["Business incubator — Univ. of Djelfa", "CCTV & camera vendors", "Retail chains & supermarkets", "Security integrators", "Cloud / hosting provider"])
bmc_cell(c(1), gy, gcw, rh2, "Key Activities",
         ["AI model development & tuning", "Software dev & integration", "Support & training"])
bmc_cell(c(1), gy + rh2 + ggap, gcw, rh2, "Key Resources",
         ["YOLOv8 models & codebase", "GPU workstation", "AI team + demo camera kit"])
bmc_cell(c(2), gy, gcw, fh2, "Value Propositions",
         ["Real-time shoplifting detection", "Works with existing CCTV", "Pose-based behaviour analysis", "Instant multi-channel alerts", "On-prem · affordable · private"])
bmc_cell(c(3), gy, gcw, rh2, "Customer Relationships",
         ["Custom zones & watch-lists", "Automated alerts & siren", "Support, training, follow-up"])
bmc_cell(c(3), gy + rh2 + ggap, gcw, rh2, "Channels",
         ["Direct sales & demos", "Security exhibitions", "Integrators + online marketing"])
bmc_cell(c(4), gy, gcw, fh2, "Customer Segments",
         ["Shops & supermarkets", "Shopping malls", "Warehouses & storage", "Security companies", "Public facilities"])
cost_w = gcw * 3 + 2 * ggap
rev_w = gcw * 2 + ggap
bmc_cell(gx, bot_y, cost_w, bot_h, "Cost Structure",
         ["Salaries & engineering", "Rent & utilities", "Hardware & cameras", "Hosting & marketing"], inline=True)
bmc_cell(gx + cost_w + ggap, bot_y, rev_w, bot_h, "Revenue Streams",
         ["Monthly subscriptions", "Installation & setup fees", "Support & training"], inline=True)
finish(s, 'fade')
notes(s, "Business Model Canvas (report Ch. 4) — 9 cells in the classic 5-column layout, THEFT DETECTION wordmark on top; cells cascade KP→Revenue.")

# ============================================================================ SLIDE CONCLUSION
s = new_slide(); header(s, "General Conclusion", "What this project proves", 0)
conclusion = [
    ("Objective achieved", "Theft Detection spots shoplifting and suspicious behaviour in real time, on a single workstation."),
    ("Behavioural intelligence", "it follows each person over time, recognising concealment, loitering and zone intrusion as they happen."),
    ("A complete product", "multi-camera monitoring, security zones, face watch-lists, evidence archive, video analysis and instant alerts."),
    ("Affordable & privacy-aware", "commodity hardware and existing cameras, with every byte of data kept on-premises."),
    ("A foundation to build on", "a dedicated theft model, edge deployment, a mobile app and advanced analytics come next."),
]
yy = Inches(2.42)
rowh = Inches(0.74)
for k, (lead, rest) in enumerate(conclusion):
    card(s, MARGIN, yy, CW, rowh, fill=CARD if k % 2 == 0 else BG2, line=LINE)
    chip(s, MARGIN + Inches(0.22), yy + int((rowh - Inches(0.40)) / 2), Inches(0.40), f'{k+1:02d}', fill=BLUE, fsize=10.5)
    text(s, MARGIN + Inches(0.92), yy, CW - Inches(1.25), rowh,
         [{'runs': [{'t': lead + '  —  ', 'size': 13, 'color': WHITE, 'bold': True},
                    {'t': rest, 'size': 12.5, 'color': MUTE}], 'ls': 1.1}], anchor=MIDDLE)
    mark()
    yy += rowh + Inches(0.12)
finish(s, 'morph')
notes(s, "General conclusion — 5 numbered card rows (lead bold + detail muted) distilled from the report's conclusion; rows cascade, Morph into the closing slide.")

# ============================================================================ SLIDE 30
s = new_slide(); grad(s, BG, GRAD2, 60)
chrome(s, 0); dotgrid(s, Inches(9.9), Inches(4.55), 9, 6)
try:
    s.shapes.add_picture(LOGO, SW - MARGIN - Inches(1.0), Inches(0.85), height=Inches(1.0))
except Exception: pass
anim_reset()
text(s, MARGIN, Inches(2.05), Inches(10), Inches(0.5), [{'t': 'THANK YOU', 'size': 14, 'color': CYAN, 'bold': True}], track=4)
mark()
text(s, MARGIN, Inches(2.6), Inches(10.5), Inches(1.7),
     [{'t': 'Let’s make every camera\nsee what matters.', 'size': 42, 'color': WHITE, 'bold': True, 'font': TITLE, 'ls': 1.1}])
mark()
rect(s, MARGIN, Inches(4.62), Inches(0.9), Pt(3), fill=CYAN); mark()
card(s, MARGIN, Inches(5.05), Inches(5.9), Inches(0.64), fill=CARD2, line=CYAN, lw=1.4, radius=0.5)
text(s, MARGIN, Inches(5.05), Inches(5.9), Inches(0.64),
     [{'runs': [{'t': 'Live demo    ', 'size': 12.5, 'color': MUTE},
                {'t': 'theft-detection-dusky.vercel.app', 'size': 14, 'color': CYAN, 'bold': True}], 'align': CENTER}],
     align=CENTER, anchor=MIDDLE)
mark()
finish(s, 'morph')
notes(s, "Closing hero — logo top-right, demo URL in a cyan pill. Morph in; elements cascade. End on the vision + URL.")

# --------------------------------------------------------------------- save + verify
out = os.environ.get("DECK_OUT", "Theft_Detection_Pitch_v2.pptx")
prs.save(out)
chk = Presentation(out)
print(f"Saved {out}  |  slides: {len(chk.slides._sldIdLst)}  |  size: {os.path.getsize(out)//1024} KB")
